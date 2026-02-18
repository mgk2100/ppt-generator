"""
ppt_utils.py — 비시각적 유틸리티
디자인을 제약하는 코드 없음. 인프라 헬퍼만 포함.
"""

import os
import shutil
import subprocess
import platform
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


BASE_DIR = Path(__file__).parent
REF_DIR = BASE_DIR / "ref"
OUTPUT_DIR = BASE_DIR / "output"
DEFAULT_TEMPLATE = REF_DIR / "표지.pptx"
FONTS_DIR = REF_DIR / "fonts"


def ensure_fonts():
    """ref/fonts/ 의 .ttf 폰트를 시스템에 설치한다. 이미 설치됐으면 스킵."""
    if not FONTS_DIR.exists():
        return False

    font_files = list(FONTS_DIR.glob("*.ttf"))
    if not font_files:
        return False

    system = platform.system()

    if system == "Linux":
        dest_dir = Path.home() / ".local" / "share" / "fonts"
    elif system == "Darwin":
        dest_dir = Path.home() / "Library" / "Fonts"
    elif system == "Windows":
        dest_dir = Path(os.environ.get("LOCALAPPDATA", "")) / "Microsoft" / "Windows" / "Fonts"
    else:
        return False

    dest_dir.mkdir(parents=True, exist_ok=True)
    installed = False

    for f in font_files:
        dest = dest_dir / f.name
        if not dest.exists():
            shutil.copy2(f, dest)
            installed = True

    if installed and system == "Linux":
        try:
            subprocess.run(["fc-cache", "-fv"], capture_output=True, check=True)
        except (subprocess.CalledProcessError, FileNotFoundError):
            pass

    return True


def load_template(page_numbers=True):
    """표지.pptx를 로드하고 샘플 슬라이드를 제거하여 빈 Presentation을 반환한다."""
    prs = Presentation(str(DEFAULT_TEMPLATE))

    # 샘플 슬라이드 제거
    slide_ids = [slide.slide_id for slide in prs.slides]
    for sid in slide_ids:
        idx = next(
            (i for i, s in enumerate(prs.slides._sldIdLst) if s.id == sid), -1
        )
        if idx >= 0:
            rId = prs.slides._sldIdLst[idx].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[idx]

    return prs


def get_layout(prs, name):
    """이름으로 슬라이드 레이아웃을 찾는다. 없으면 ValueError."""
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == name:
            return layout
    raise ValueError(f"레이아웃 '{name}'을 찾을 수 없습니다.")


def clear_placeholders(slide, keep=None):
    """마스터 슬라이드에서 상속된 유령 플레이스홀더/텍스트를 제거한다.

    Args:
        slide: 슬라이드 객체
        keep: 유지할 플레이스홀더 idx 리스트
    """
    if keep is None:
        keep = []

    ghost_texts = [
        "마스터 텍스트 스타일 편집",
        "마스터 텍스트 스타일을 편집합니다",
        "마스터 제목 스타일 편집",
        "제목을 추가하려면 클릭하십시오",
        "제목을 입력하십시오",
        "부제목을 입력하십시오",
        "텍스트를 입력하십시오",
        "내용을 입력하십시오",
        "텍스트를 추가하려면 클릭하십시오",
        "Click to edit Master text styles",
        "Click to edit Master title style",
        "Click to add title",
        "Click to add text",
        "Click to add subtitle",
    ]

    to_remove = []

    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx in keep:
            continue
        if ph.has_text_frame:
            text = ph.text_frame.text.strip().rstrip(".")
            if not text or any(g in text or text in g for g in ghost_texts):
                to_remove.append(ph)

    for shape in slide.shapes:
        if shape in to_remove:
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip().rstrip(".")
            if any(g in text or text in g for g in ghost_texts):
                to_remove.append(shape)

    for shape in to_remove:
        shape._element.getparent().remove(shape._element)


def set_cell_anchor(cell, anchor="ctr"):
    """테이블 셀 세로정렬 XML 워크어라운드.

    Args:
        cell: python-pptx 테이블 셀
        anchor: 't' (위), 'ctr' (가운데), 'b' (아래)
    """
    from lxml import etree

    a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    tc = cell._tc

    # tcPr 에 anchor 설정
    tcPr = next((c for c in tc if "tcPr" in c.tag), None)
    if tcPr is None:
        tcPr = etree.Element(f"{a_ns}tcPr")
        tc.insert(0, tcPr)
    tcPr.set("anchor", anchor)

    # txBody > bodyPr 에도 anchor 설정
    txBody = next((c for c in tc if "txBody" in c.tag), None)
    if txBody is None:
        _ = cell.text_frame  # txBody 생성
        txBody = next((c for c in tc if "txBody" in c.tag), None)

    if txBody is not None:
        bodyPr = next((c for c in txBody if "bodyPr" in c.tag), None)
        if bodyPr is None:
            bodyPr = etree.Element(f"{a_ns}bodyPr")
            txBody.insert(0, bodyPr)
        bodyPr.set("anchor", anchor)


def add_arrowhead(connector):
    """커넥터에 화살표 머리를 추가한다 (python-pptx에 네이티브 API 없음)."""
    connector.line._ln.append(
        connector.line._ln.makeelement(
            qn("a:tailEnd"),
            {"type": "triangle", "w": "med", "len": "med"},
        )
    )


# ---------------------------------------------------------------------------
# 신규 유틸리티 — XML 배관 코드 캡슐화 (시각적 의견 없음)
# ---------------------------------------------------------------------------


def add_shadow(shape, blur_pt=4, dist_pt=3, direction=2700000,
               opacity_pct=40, color=None):
    """도형에 outerShadow를 추가한다.

    Args:
        shape: python-pptx 도형 객체
        blur_pt: 블러 반경 (포인트 단위)
        dist_pt: 그림자 거리 (포인트 단위)
        direction: 그림자 방향 (EMU 각도, 기본 270° = 아래)
        opacity_pct: 그림자 불투명도 (0-100)
        color: RGBColor 또는 (r,g,b) 튜플. None이면 검정
    """
    if color is None:
        r, g, b = 0, 0, 0
    elif isinstance(color, RGBColor):
        r, g, b = color[0], color[1], color[2]
    else:
        r, g, b = color

    alpha_val = int(opacity_pct * 1000)  # 40% → 40000
    blur_emu = str(Pt(blur_pt))
    dist_emu = str(Pt(dist_pt))
    hex_color = f"{r:02X}{g:02X}{b:02X}"

    spPr = shape._element.spPr if hasattr(shape._element, 'spPr') else None
    if spPr is None:
        spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return

    # effectLst 찾기/생성
    effectLst = spPr.find(qn("a:effectLst"))
    if effectLst is None:
        effectLst = spPr.makeelement(qn("a:effectLst"), {})
        spPr.append(effectLst)

    # 기존 outerShdw 제거
    for old in effectLst.findall(qn("a:outerShdw")):
        effectLst.remove(old)

    outerShdw = effectLst.makeelement(qn("a:outerShdw"), {
        "blurRad": blur_emu,
        "dist": dist_emu,
        "dir": str(direction),
        "rotWithShape": "0",
    })
    srgbClr = outerShdw.makeelement(qn("a:srgbClr"), {"val": hex_color})
    alphaElem = srgbClr.makeelement(qn("a:alpha"), {"val": str(alpha_val)})
    srgbClr.append(alphaElem)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)


def set_shape_opacity(shape, opacity_pct):
    """도형의 채우기 투명도를 설정한다.

    Args:
        shape: python-pptx 도형 객체 (solidFill이 이미 적용된 상태여야 함)
        opacity_pct: 불투명도 (0=완전투명, 100=불투명)
    """
    alpha_val = str(int(opacity_pct * 1000))  # 50% → 50000

    spPr = shape._element.spPr if hasattr(shape._element, 'spPr') else None
    if spPr is None:
        spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return

    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is None:
        return

    # srgbClr 또는 schemeClr 찾기
    color_elem = solidFill.find(qn("a:srgbClr"))
    if color_elem is None:
        color_elem = solidFill.find(qn("a:schemeClr"))
    if color_elem is None:
        return

    # 기존 alpha 제거 후 새로 추가
    for old in color_elem.findall(qn("a:alpha")):
        color_elem.remove(old)
    alpha_elem = color_elem.makeelement(qn("a:alpha"), {"val": alpha_val})
    color_elem.append(alpha_elem)


def add_gradient_stop(shape, position, r, g, b):
    """그라디언트 fill에 추가 stop을 삽입한다.

    shape.fill.gradient() 로 기본 2-stop 설정 후,
    3번째 이상의 stop을 추가할 때 사용한다.

    Args:
        shape: python-pptx 도형 (이미 gradient fill 적용된 상태)
        position: 0.0~1.0 (0=시작, 1=끝)
        r, g, b: 정수 0-255
    """
    spPr = shape._element.spPr if hasattr(shape._element, 'spPr') else None
    if spPr is None:
        spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return

    gradFill = spPr.find(qn("a:gradFill"))
    if gradFill is None:
        return

    gsLst = gradFill.find(qn("a:gsLst"))
    if gsLst is None:
        gsLst = gradFill.makeelement(qn("a:gsLst"), {})
        gradFill.insert(0, gsLst)

    pos_val = str(int(position * 100000))  # 0.5 → 50000
    hex_color = f"{r:02X}{g:02X}{b:02X}"

    gs = gsLst.makeelement(qn("a:gs"), {"pos": pos_val})
    srgbClr = gs.makeelement(qn("a:srgbClr"), {"val": hex_color})
    gs.append(srgbClr)
    gsLst.append(gs)


def make_icon_circle(slide, x, y, size, fill_color, text="",
                     font_size=10, font_color=None):
    """원형 아이콘/배지를 생성한다 (OVAL + fill + 중앙정렬 텍스트).

    번호 배지, 상태 표시, 아이콘 대체 등에 사용.
    색상/크기는 파라미터이므로 시각적 의견 없음.

    Args:
        slide: 슬라이드 객체
        x, y: 위치 (Inches/Emu)
        size: 원 지름 (Inches/Emu)
        fill_color: RGBColor
        text: 원 안에 들어갈 텍스트
        font_size: 포인트 단위
        font_color: RGBColor. None이면 brightness_check로 자동 결정

    Returns:
        생성된 도형 객체
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # 테두리 없음

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(font_size)
        run.font.bold = True

        if font_color is None:
            is_bright = brightness_check(fill_color[0], fill_color[1], fill_color[2])
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33) if is_bright else RGBColor(0xFF, 0xFF, 0xFF)
        else:
            run.font.color.rgb = font_color

        # 세로 중앙정렬
        tf_body = shape.text_frame._txBody
        bodyPr = tf_body.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("anchor", "ctr")

    return shape


def brightness_check(r, g, b):
    """배경 색상의 밝기를 판단한다.

    Args:
        r, g, b: 정수 0-255

    Returns:
        True = 밝은 배경 (어두운 텍스트 사용)
        False = 어두운 배경 (흰색 텍스트 사용)
    """
    return (r * 0.299 + g * 0.587 + b * 0.114) > 160


# ---------------------------------------------------------------------------
# 텍스트 편의 함수 — 반복 보일러플레이트 제거 (시각적 의견 없음)
# ---------------------------------------------------------------------------


def add_textbox(slide, x, y, w, h, text, font_name=None, font_size=12,
                color=None, bold=False, align=PP_ALIGN.LEFT, word_wrap=True):
    """텍스트박스를 추가하고 단일 단락을 설정한다.

    Args:
        slide: 슬라이드 객체
        x, y, w, h: 위치/크기 (Inches/Emu)
        text: 텍스트 내용
        font_name: 폰트 이름 (None이면 기본)
        font_size: 포인트 단위
        color: RGBColor (None이면 기본)
        bold: 볼드 여부
        align: PP_ALIGN 정렬
        word_wrap: 자동 줄바꿈

    Returns:
        생성된 텍스트박스 도형 객체
    """
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    if color:
        run.font.color.rgb = color

    return txBox


def add_para(text_frame, text, font_name=None, font_size=12, color=None,
             bold=False, align=PP_ALIGN.LEFT, space_before=None,
             space_after=None):
    """기존 text_frame에 새 단락을 추가한다.

    Args:
        text_frame: python-pptx TextFrame 객체
        text: 텍스트 내용
        font_name: 폰트 이름 (None이면 기본)
        font_size: 포인트 단위
        color: RGBColor (None이면 기본)
        bold: 볼드 여부
        align: PP_ALIGN 정렬
        space_before: 단락 전 간격 (Pt 단위 값)
        space_after: 단락 후 간격 (Pt 단위 값)

    Returns:
        생성된 단락 객체
    """
    p = text_frame.add_paragraph()
    p.alignment = align
    if space_before is not None:
        p.space_before = space_before
    if space_after is not None:
        p.space_after = space_after

    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if font_name:
        run.font.name = font_name
    if color:
        run.font.color.rgb = color

    return p


def set_body_anchor(shape, anchor="ctr"):
    """도형 텍스트 프레임의 세로정렬을 설정한다.

    set_cell_anchor의 형제 함수. 테이블 셀이 아닌 일반 도형용.

    Args:
        shape: python-pptx 도형 객체 (text_frame을 가진)
        anchor: 't' (위), 'ctr' (가운데), 'b' (아래)
    """
    if not shape.has_text_frame:
        return
    bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        bodyPr.set("anchor", anchor)
