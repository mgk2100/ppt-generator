#!/usr/bin/env python3
"""
PPT Generator - 템플릿 기반 PowerPoint 생성기
표지.pptx 템플릿을 기반으로 보고서 PPT를 자동 생성합니다.
"""

import argparse
import json
import yaml
import os
import sys
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData, ChartData
from pptx.oxml.ns import qn
import math
import copy
import shutil
import subprocess
import platform


# 기본 설정
BASE_DIR = Path(__file__).parent
TEMPLATE_DIR = BASE_DIR / "templates"
OUTPUT_DIR = BASE_DIR / "output"
DEFAULT_TEMPLATE = TEMPLATE_DIR / "표지.pptx"
FONTS_DIR = BASE_DIR / "fonts"


def install_fonts():
    """현대하모니 폰트를 시스템에 설치합니다."""
    if not FONTS_DIR.exists():
        print(f"폰트 디렉토리가 없습니다: {FONTS_DIR}")
        return False

    # 설치할 폰트 파일 목록
    font_files = list(FONTS_DIR.glob("*.ttf"))
    if not font_files:
        print("설치할 폰트 파일이 없습니다.")
        return False

    system = platform.system()

    if system == "Linux":
        # Linux: ~/.local/share/fonts에 복사
        user_fonts_dir = Path.home() / ".local" / "share" / "fonts"
        user_fonts_dir.mkdir(parents=True, exist_ok=True)

        installed = False
        for font_file in font_files:
            dest = user_fonts_dir / font_file.name
            if not dest.exists():
                shutil.copy2(font_file, dest)
                print(f"폰트 설치: {font_file.name}")
                installed = True

        if installed:
            # 폰트 캐시 업데이트
            try:
                subprocess.run(["fc-cache", "-fv"], capture_output=True, check=True)
                print("폰트 캐시 업데이트 완료")
            except (subprocess.CalledProcessError, FileNotFoundError):
                print("폰트 캐시 업데이트 실패 (fc-cache 명령어를 찾을 수 없음)")
        else:
            print("현대하모니 폰트가 이미 설치되어 있습니다.")
        return True

    elif system == "Darwin":
        # macOS: ~/Library/Fonts에 복사
        user_fonts_dir = Path.home() / "Library" / "Fonts"
        user_fonts_dir.mkdir(parents=True, exist_ok=True)

        for font_file in font_files:
            dest = user_fonts_dir / font_file.name
            if not dest.exists():
                shutil.copy2(font_file, dest)
                print(f"폰트 설치: {font_file.name}")
        print("macOS 폰트 설치 완료")
        return True

    elif system == "Windows":
        # Windows: %LOCALAPPDATA%\Microsoft\Windows\Fonts에 복사
        user_fonts_dir = Path(os.environ.get("LOCALAPPDATA", "")) / "Microsoft" / "Windows" / "Fonts"
        user_fonts_dir.mkdir(parents=True, exist_ok=True)

        for font_file in font_files:
            dest = user_fonts_dir / font_file.name
            if not dest.exists():
                shutil.copy2(font_file, dest)
                print(f"폰트 설치: {font_file.name}")
        print("Windows 폰트 설치 완료")
        return True

    else:
        print(f"지원하지 않는 운영체제: {system}")
        return False


def check_fonts_installed():
    """현대하모니 폰트가 시스템에 설치되어 있는지 확인합니다."""
    system = platform.system()

    if system == "Linux":
        user_fonts_dir = Path.home() / ".local" / "share" / "fonts"
        required_fonts = ["현대하모니 M.ttf", "현대하모니 L.ttf"]
        for font in required_fonts:
            if not (user_fonts_dir / font).exists():
                return False
        return True
    elif system == "Darwin":
        user_fonts_dir = Path.home() / "Library" / "Fonts"
        required_fonts = ["현대하모니 M.ttf", "현대하모니 L.ttf"]
        for font in required_fonts:
            if not (user_fonts_dir / font).exists():
                return False
        return True
    elif system == "Windows":
        user_fonts_dir = Path(os.environ.get("LOCALAPPDATA", "")) / "Microsoft" / "Windows" / "Fonts"
        required_fonts = ["현대하모니 M.ttf", "현대하모니 L.ttf"]
        for font in required_fonts:
            if not (user_fonts_dir / font).exists():
                return False
        return True
    return False


# 프로그램 시작 시 폰트 자동 설치
if not check_fonts_installed():
    print("현대하모니 폰트를 설치합니다...")
    install_fonts()


class DesignSystem:
    """통합 디자인 시스템 - 일관된 스타일 관리"""

    # 기본 회사 브랜드 색상 (ref_1.pptx 기반 개선)
    DEFAULT_BRAND_COLORS = {
        "primary": (40, 55, 78),        # #28374E - 다크 네이비 (ref_1 메인)
        "secondary": (79, 129, 189),    # #4F81BD - 미드 블루
        "accent": (31, 73, 125),        # #1F497D - 딥 블루
        "highlight": (255, 192, 0),     # 골드/노랑
        "success": (53, 162, 159),      # #35A29F - 티일 그린
        "warning": (255, 167, 109),     # #FFA76D - 소프트 오렌지
        "danger": (237, 102, 102),      # #ED6666 - 소프트 레드
        "light": (232, 237, 244),       # #E8EDF4 - 콘텐츠 박스 배경
        "dark": (51, 51, 51),           # #333333 - 메인 텍스트
        "text": (51, 51, 51),           # #333333 - 본문 텍스트
        "white": (255, 255, 255),
        "black": (0, 0, 0),
        # 추가 색상 (ref_1.pptx 기반)
        "content_box": (232, 237, 244), # #E8EDF4 - 콘텐츠 박스 배경
        "header_bg": (40, 55, 78),      # #28374E - 헤더 배경
        "card_border": (220, 220, 220), # 카드 테두리
        "teal": (11, 102, 105),         # #0B6669 - 다크 티일
        "navy": (8, 24, 83),            # #081853 - 딥 네이비
    }

    # 기본 그라데이션 팔레트 (ref_1.pptx 기반)
    DEFAULT_GRADIENT = [
        (40, 55, 78),     # #28374E - 가장 진함
        (31, 73, 125),    # #1F497D
        (79, 129, 189),   # #4F81BD
        (126, 155, 200),  # 중간
        (181, 211, 235),  # #B5D3EB - 가장 연함
    ]

    # 기본 폰트 설정 (현대하모니 폰트)
    # 현대하모니M: 제목용 (Medium)
    # 현대하모니L: 본문용 (Light)
    FONT_TITLE = "현대하모니 M"      # 대주제용
    FONT_BODY = "현대하모니 L"       # 본문/내용용

    DEFAULT_FONTS = {
        "cover_title": {"name": "현대하모니 M", "size": 44, "bold": True},      # 표지 주제 (ref_1: 44pt)
        "cover_date": {"name": "현대하모니 M", "size": 14, "bold": False},      # 표지 날짜
        "cover_author": {"name": "현대하모니 L", "size": 14, "bold": True},     # 표지 작성자
        "cover_type": {"name": "현대하모니 L", "size": 12, "bold": False},      # 표지 보고유형
        "title": {"name": "현대하모니 M", "size": 20, "bold": True},            # 슬라이드 대주제 (ref_1: 20pt)
        "section": {"name": "현대하모니 M", "size": 14, "bold": True},          # 섹션 헤더 (● 불릿)
        "subtitle": {"name": "현대하모니 L", "size": 16, "bold": True},
        "heading": {"name": "현대하모니 L", "size": 14, "bold": True},
        "subheading": {"name": "현대하모니 L", "size": 12, "bold": True},
        "body": {"name": "현대하모니 L", "size": 12, "bold": False},            # 본문 (ref_1: 11-12pt)
        "caption": {"name": "현대하모니 L", "size": 11, "bold": False},
        "small": {"name": "현대하모니 L", "size": 9, "bold": False},            # 작은 텍스트 (ref_1: 9pt)
    }

    # 기본 슬라이드 레이아웃 설정 (인치)
    DEFAULT_LAYOUT = {
        "margin_left": 0.4,
        "margin_right": 0.4,
        "margin_top": 0.9,
        "margin_bottom": 0.5,
        "content_width": 10.0,
        "title_height": 0.5,
        "spacing": 0.15,
    }

    # 카드 스타일 옵션 (9가지)
    # classic: [백업] 기존 스타일 (좌측 컬러바 + 상단 원형 아이콘)
    # gradient: 상단 그라데이션 헤더 + 아이콘
    # modern: 좌측 큰 아이콘 강조형
    # solid: 전체 컬러 카드
    # outline: 테두리 강조 + 상단 아이콘
    # minimal: 미니멀 - 하단 컬러 라인만
    # banner: 배너 스타일 - 상단 풀 컬러 배너
    # split: 분할 카드 - 상단 컬러/하단 화이트
    # accent: 좌측 두꺼운 악센트 바 + 큰 아이콘
    CARD_STYLES = ["classic", "gradient", "modern", "solid", "outline", "minimal", "banner", "split", "accent"]
    DEFAULT_CARD_STYLE = "gradient"  # 새 기본값

    # 사전 정의된 테마
    PRESET_THEMES = {
        "default": {},  # 기본값 사용
        "dark": {
            "colors": {
                "primary": (33, 37, 41),
                "secondary": (52, 58, 64),
                "accent": (0, 123, 255),
                "light": (73, 80, 87),
                "text": (248, 249, 250),
            }
        },
        "green": {
            "colors": {
                "primary": (25, 135, 84),
                "secondary": (32, 201, 151),
                "accent": (13, 110, 253),
            }
        },
        "purple": {
            "colors": {
                "primary": (111, 66, 193),
                "secondary": (214, 51, 132),
                "accent": (102, 16, 242),
            }
        },
        "warm": {
            "colors": {
                "primary": (220, 53, 69),
                "secondary": (253, 126, 20),
                "accent": (255, 193, 7),
                "success": (25, 135, 84),
            }
        },
    }

    def __init__(self, theme_path: str = None, theme_name: str = None):
        """디자인 시스템 초기화

        Args:
            theme_path: 외부 테마 파일 경로 (YAML/JSON)
            theme_name: 사전 정의된 테마 이름 (default, dark, green, purple, warm)
        """
        # 기본값으로 초기화
        self._init_defaults()

        # 사전 정의된 테마 적용
        if theme_name and theme_name in self.PRESET_THEMES:
            self._apply_theme(self.PRESET_THEMES[theme_name])

        # 외부 테마 파일 로드
        if theme_path:
            self.load_theme(theme_path)

    def _init_defaults(self):
        """기본값으로 초기화"""
        # 색상을 RGBColor 객체로 변환
        self.BRAND_COLORS = {
            k: RGBColor(*v) for k, v in self.DEFAULT_BRAND_COLORS.items()
        }
        self.GRADIENT_BLUE = [RGBColor(*c) for c in self.DEFAULT_GRADIENT]
        self.FONTS = copy.deepcopy(self.DEFAULT_FONTS)
        self.LAYOUT = copy.deepcopy(self.DEFAULT_LAYOUT)
        self.card_style = self.DEFAULT_CARD_STYLE

    def _apply_theme(self, theme_config: Dict[str, Any]):
        """테마 설정 적용"""
        # 색상 적용
        if "colors" in theme_config:
            for name, color in theme_config["colors"].items():
                if isinstance(color, (list, tuple)) and len(color) == 3:
                    self.BRAND_COLORS[name] = RGBColor(*color)
                elif isinstance(color, str) and color.startswith("#"):
                    # HEX 색상 지원
                    hex_color = color.lstrip("#")
                    r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                    self.BRAND_COLORS[name] = RGBColor(r, g, b)

        # 그라데이션 적용
        if "gradient" in theme_config:
            self.GRADIENT_BLUE = []
            for color in theme_config["gradient"]:
                if isinstance(color, (list, tuple)):
                    self.GRADIENT_BLUE.append(RGBColor(*color))
                elif isinstance(color, str) and color.startswith("#"):
                    hex_color = color.lstrip("#")
                    r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                    self.GRADIENT_BLUE.append(RGBColor(r, g, b))

        # 폰트 적용
        if "fonts" in theme_config:
            for style, settings in theme_config["fonts"].items():
                if style in self.FONTS:
                    self.FONTS[style].update(settings)
                else:
                    self.FONTS[style] = settings

        # 레이아웃 적용
        if "layout" in theme_config:
            self.LAYOUT.update(theme_config["layout"])

        # 카드 스타일 적용
        if "card_style" in theme_config:
            if theme_config["card_style"] in self.CARD_STYLES:
                self.card_style = theme_config["card_style"]

    def load_theme(self, theme_path: str) -> bool:
        """외부 테마 파일 로드

        Args:
            theme_path: 테마 파일 경로 (YAML 또는 JSON)

        Returns:
            로드 성공 여부
        """
        theme_file = Path(theme_path)
        if not theme_file.exists():
            print(f"경고: 테마 파일을 찾을 수 없습니다: {theme_path}")
            return False

        try:
            if theme_file.suffix in [".yaml", ".yml"]:
                with open(theme_file, "r", encoding="utf-8") as f:
                    theme_config = yaml.safe_load(f)
            else:
                with open(theme_file, "r", encoding="utf-8") as f:
                    theme_config = json.load(f)

            self._apply_theme(theme_config)
            return True

        except Exception as e:
            print(f"테마 로드 오류: {e}")
            return False

    def _rgb_to_list(self, color: RGBColor) -> List[int]:
        """RGBColor를 [r, g, b] 리스트로 변환"""
        # RGBColor는 str()로 hex 문자열을 반환 (예: "003366")
        hex_str = str(color)
        r = int(hex_str[0:2], 16)
        g = int(hex_str[2:4], 16)
        b = int(hex_str[4:6], 16)
        return [r, g, b]

    def save_theme(self, output_path: str):
        """현재 테마를 파일로 저장

        Args:
            output_path: 저장할 파일 경로
        """
        theme_config = {
            "colors": {
                name: self._rgb_to_list(color)
                for name, color in self.BRAND_COLORS.items()
            },
            "gradient": [
                self._rgb_to_list(c) for c in self.GRADIENT_BLUE
            ],
            "fonts": self.FONTS,
            "layout": self.LAYOUT,
        }

        output_file = Path(output_path)
        with open(output_file, "w", encoding="utf-8") as f:
            if output_file.suffix in [".yaml", ".yml"]:
                yaml.dump(theme_config, f, allow_unicode=True, default_flow_style=False)
            else:
                json.dump(theme_config, f, ensure_ascii=False, indent=2)

    def set_color(self, name: str, r: int, g: int, b: int):
        """색상 설정"""
        self.BRAND_COLORS[name] = RGBColor(r, g, b)

    def set_font(self, style: str, name: str = None, size: int = None, bold: bool = None):
        """폰트 설정"""
        if style not in self.FONTS:
            self.FONTS[style] = {"name": "맑은 고딕", "size": 14, "bold": False}
        if name:
            self.FONTS[style]["name"] = name
        if size:
            self.FONTS[style]["size"] = size
        if bold is not None:
            self.FONTS[style]["bold"] = bold


class PPTGenerator:
    """PPT 생성기 클래스"""

    # 보고 유형 매핑
    REPORT_TYPES = {
        "의사결정": "■ 의사결정    □ 보고    □ 정보공유",
        "보고": "□ 의사결정    ■ 보고    □ 정보공유",
        "정보공유": "□ 의사결정    □ 보고    ■ 정보공유",
    }

    def __init__(
        self,
        template_path: Optional[Path] = None,
        show_page_numbers: bool = True,
        theme_path: str = None,
        theme_name: str = None
    ):
        """PPT 생성기 초기화

        Args:
            template_path: 템플릿 파일 경로
            show_page_numbers: 페이지 번호 표시 여부 (기본: True)
            theme_path: 외부 테마 파일 경로 (YAML/JSON)
            theme_name: 사전 정의된 테마 이름 (default, dark, green, purple, warm)
        """
        self.template_path = template_path or DEFAULT_TEMPLATE
        self.prs = Presentation(str(self.template_path))
        self.design = DesignSystem(theme_path=theme_path, theme_name=theme_name)
        self.show_page_numbers = show_page_numbers
        self._remove_sample_slides()

    def _get_content_layout_name(self) -> str:
        """페이지 번호 설정에 따른 레이아웃 이름 반환"""
        if self.show_page_numbers:
            return "제목 및 내용"
        else:
            return "제목 및 내용 (페이지 번호 삭제)"

    def _remove_sample_slides(self):
        """템플릿의 샘플 슬라이드 제거"""
        slide_ids = [slide.slide_id for slide in self.prs.slides]
        for slide_id in slide_ids:
            rId = self.prs.slides._sldIdLst[self._get_slide_index(slide_id)].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[self._get_slide_index(slide_id)]

    def _get_slide_index(self, slide_id: int) -> int:
        for idx, sldId in enumerate(self.prs.slides._sldIdLst):
            if sldId.id == slide_id:
                return idx
        return -1

    def _get_layout(self, layout_name: str):
        for layout in self.prs.slide_masters[0].slide_layouts:
            if layout.name == layout_name:
                return layout
        raise ValueError(f"레이아웃 '{layout_name}'을 찾을 수 없습니다.")

    def _clear_unused_placeholders(self, slide, used_placeholder_idx: list = None):
        """슬라이드에서 사용하지 않는 플레이스홀더 제거

        마스터 슬라이드에서 상속된 기본 플레이스홀더 텍스트
        (예: '마스터 텍스트 스타일 편집')를 제거합니다.

        Args:
            slide: 슬라이드 객체
            used_placeholder_idx: 사용 중인 플레이스홀더 idx 리스트
        """
        if used_placeholder_idx is None:
            used_placeholder_idx = []

        # 제거할 기본 플레이스홀더 텍스트 패턴
        default_texts = [
            "마스터 텍스트 스타일 편집",
            "마스터 텍스트 스타일을 편집합니다",
            "마스터 제목 스타일 편집",
            "제목을 추가하려면 클릭하십시오",
            "제목을 입력하십시오",
            "부제목을 입력하십시오",
            "텍스트를 입력하십시오",
            "내용을 입력하십시오",
            "텍스트를 추가하려면 클릭하십시오",
            "Click to edit Master text styles",
            "Click to edit Master title style",
            "Click to add title",
            "Click to add text",
            "Click to add subtitle",
        ]

        shapes_to_remove = []

        # 먼저 placeholders 컬렉션에서 사용하지 않는 것들 제거
        for placeholder in list(slide.placeholders):
            ph_idx = placeholder.placeholder_format.idx
            # 사용 중인 플레이스홀더는 건너뜀
            if ph_idx in used_placeholder_idx:
                continue

            if placeholder.has_text_frame:
                text = placeholder.text_frame.text.strip()
                # 빈 텍스트이거나 기본 플레이스홀더 텍스트인 경우
                if not text:
                    shapes_to_remove.append(placeholder)
                else:
                    text_clean = text.rstrip('.')
                    for default_text in default_texts:
                        if default_text in text_clean or text_clean in default_text:
                            shapes_to_remove.append(placeholder)
                            break

        # 일반 shapes에서도 기본 텍스트 패턴 체크
        for shape in slide.shapes:
            if shape in shapes_to_remove:
                continue

            if shape.has_text_frame:
                text = shape.text_frame.text.strip().rstrip('.')

                # 기본 플레이스홀더 텍스트인지 확인
                for default_text in default_texts:
                    if default_text in text or text in default_text:
                        if shape not in shapes_to_remove:
                            shapes_to_remove.append(shape)
                        break

        # 플레이스홀더 제거 (shapes 컬렉션에서 제거)
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    def _apply_text_style(self, paragraph, style_name: str, color: RGBColor = None):
        """텍스트 스타일 적용"""
        style = self.design.FONTS.get(style_name, self.design.FONTS["body"])
        paragraph.font.name = style["name"]
        paragraph.font.size = Pt(style["size"])
        paragraph.font.bold = style["bold"]
        if color:
            paragraph.font.color.rgb = color

    def _add_decorative_line(self, slide, y: float, color: RGBColor = None, width: float = None):
        """장식 구분선 추가"""
        line_color = color or self.design.BRAND_COLORS["primary"]
        line_width = width or self.design.LAYOUT["content_width"]

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(self.design.LAYOUT["margin_left"]),
            Inches(y),
            Inches(line_width),
            Inches(0.03)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = line_color
        shape.line.fill.background()
        return shape

    def _add_accent_bar(self, slide, x: float, y: float, height: float, color: RGBColor = None):
        """강조 세로 바 추가"""
        bar_color = color or self.design.BRAND_COLORS["accent"]
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(0.08),
            Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bar_color
        shape.line.fill.background()
        return shape

    def _add_icon_box(
        self,
        slide,
        icon_text: str,
        x: float,
        y: float,
        size: float = 0.5,
        bg_color: RGBColor = None,
        text_color: RGBColor = None
    ):
        """아이콘 스타일 박스 추가 (숫자나 기호용)"""
        bg = bg_color or self.design.BRAND_COLORS["primary"]
        fg = text_color or self.design.BRAND_COLORS["white"]

        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x), Inches(y),
            Inches(size), Inches(size)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = False
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(icon_text)
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(int(size * 20))
        p.font.bold = True
        p.font.color.rgb = fg

        return shape

    def _add_shadow_box(
        self,
        slide,
        x: float,
        y: float,
        width: float,
        height: float,
        fill_color: RGBColor = None,
        border_color: RGBColor = None,
        shadow_offset: float = 0.05,
        shadow_color: RGBColor = None,
        corner_radius: bool = True
    ):
        """그림자 효과가 있는 박스 추가

        Args:
            slide: 슬라이드 객체
            x, y: 박스 위치 (인치)
            width, height: 박스 크기 (인치)
            fill_color: 박스 배경색 (기본: 흰색)
            border_color: 테두리 색상 (기본: 연한 회색)
            shadow_offset: 그림자 오프셋 (인치)
            shadow_color: 그림자 색상 (기본: 연한 회색)
            corner_radius: 둥근 모서리 사용 여부
        """
        fill = fill_color or self.design.BRAND_COLORS["white"]
        border = border_color or RGBColor(220, 220, 220)
        shadow_clr = shadow_color or RGBColor(200, 200, 200)
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE

        # 그림자 (오프셋된 회색 박스)
        shadow = slide.shapes.add_shape(
            shape_type,
            Inches(x + shadow_offset), Inches(y + shadow_offset),
            Inches(width), Inches(height)
        )
        shadow.fill.solid()
        shadow.fill.fore_color.rgb = shadow_clr
        shadow.line.fill.background()

        # 메인 박스
        box = slide.shapes.add_shape(
            shape_type,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = border
        box.line.width = Pt(1)

        return box

    # 아이콘 매핑 (공통 사용)
    ICON_MAPPING = {
        "document": "📄", "ai": "🤖", "flow": "⚙️", "server": "🖥️",
        "database": "🗄️", "link": "🔗", "settings": "⚙️", "chart": "📊",
        "code": "💻", "cloud": "☁️", "security": "🔒", "network": "🌐",
        "user": "👤", "api": "🔌", "data": "📁", "check": "✓",
        "star": "★", "heart": "♥", "lightning": "⚡", "target": "◎",
    }

    def _get_icon_text(self, icon: str, card_index: int) -> str:
        """아이콘 텍스트 변환"""
        if icon:
            return self.ICON_MAPPING.get(icon.lower(), icon[:2].upper() if len(icon) > 2 else icon)
        return str(card_index + 1)

    def _add_section_header(
        self, slide, title: str, x: float, y: float, width: float = 4.0
    ):
        """섹션 헤더 추가 (● 불릿 스타일, ref_1.pptx 기반)

        Args:
            slide: 슬라이드 객체
            title: 섹션 제목
            x, y: 위치 (인치)
            width: 너비 (인치)

        Returns:
            텍스트박스 높이 (인치)
        """
        header_height = 0.5

        # 섹션 헤더 배경 (선택적 - 좌측 컬러 라인)
        line_width = 0.08
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y + 0.05),
            Inches(line_width), Inches(header_height - 0.1)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.design.BRAND_COLORS["primary"]
        line.line.fill.background()

        # 섹션 제목 텍스트 (● 불릿 포함)
        text_box = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y),
            Inches(width), Inches(header_height)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"● {title}"
        p.font.name = self.design.FONTS["section"]["name"]
        p.font.size = Pt(self.design.FONTS["section"]["size"])
        p.font.bold = self.design.FONTS["section"]["bold"]
        p.font.color.rgb = self.design.BRAND_COLORS["primary"]

        return header_height

    def _add_content_box(
        self, slide, x: float, y: float, width: float, height: float,
        title: str = None, show_border: bool = True
    ):
        """콘텐츠 박스 배경 추가 (ref_1.pptx #E8EDF4 스타일)

        Args:
            slide: 슬라이드 객체
            x, y: 위치 (인치)
            width, height: 크기 (인치)
            title: 박스 상단 제목 (선택)
            show_border: 테두리 표시 여부

        Returns:
            (content_y, content_height): 내부 콘텐츠 시작 y좌표와 높이
        """
        title_height = 0.4 if title else 0

        # 메인 박스 배경
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = self.design.BRAND_COLORS["content_box"]

        if show_border:
            box.line.color.rgb = RGBColor(200, 210, 225)
            box.line.width = Pt(1)
        else:
            box.line.fill.background()

        # 제목이 있는 경우 헤더 영역
        if title:
            header = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y),
                Inches(width), Inches(title_height)
            )
            header.fill.solid()
            header.fill.fore_color.rgb = self.design.BRAND_COLORS["primary"]
            header.line.fill.background()

            title_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(y + 0.05),
                Inches(width - 0.3), Inches(title_height - 0.1)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.design.BRAND_COLORS["white"]

        return (y + title_height + 0.1, height - title_height - 0.15)

    def _lighten_color(self, color: RGBColor, factor: float = 0.7) -> RGBColor:
        """색상을 밝게 만듦

        Args:
            color: 원본 색상
            factor: 밝기 증가 비율 (0-1, 높을수록 밝음)

        Returns:
            밝아진 RGBColor
        """
        r = int(color[0] + (255 - color[0]) * factor)
        g = int(color[1] + (255 - color[1]) * factor)
        b = int(color[2] + (255 - color[2]) * factor)
        return RGBColor(min(255, r), min(255, g), min(255, b))

    def _add_card(
        self,
        slide,
        title: str,
        content: str,
        x: float,
        y: float,
        width: float,
        height: float,
        accent_color: RGBColor = None,
        show_shadow: bool = True,
        icon: str = None,
        card_index: int = 0,
        card_style: str = None
    ):
        """카드 스타일 컴포넌트 추가 - 스타일 디스패처

        Args:
            card_style: 카드 스타일 (9가지)
                - classic: [백업] 기존 디자인
                - gradient: 상단 그라데이션 헤더
                - modern: 좌측 큰 아이콘
                - solid: 전체 컬러
                - outline: 테두리 강조
                - minimal: 미니멀
                - banner: 배너 스타일
                - split: 분할 카드
                - accent: 악센트 바 강조
        """
        style = card_style or self.design.card_style

        style_map = {
            "classic": self._add_card_classic,
            "gradient": self._add_card_gradient,
            "modern": self._add_card_modern,
            "solid": self._add_card_solid,
            "outline": self._add_card_outline,
            "minimal": self._add_card_minimal,
            "banner": self._add_card_banner,
            "split": self._add_card_split,
            "accent": self._add_card_accent,
        }

        func = style_map.get(style, self._add_card_gradient)
        return func(slide, title, content, x, y, width, height,
                   accent_color, show_shadow, icon, card_index)

    def _add_card_classic(
        self, slide, title: str, content: str, x: float, y: float,
        width: float, height: float, accent_color: RGBColor = None,
        show_shadow: bool = True, icon: str = None, card_index: int = 0
    ):
        """[백업] 기존 클래식 카드 스타일 - 좌측 컬러바 + 상단 원형 아이콘"""
        accent = accent_color or self.design.BRAND_COLORS["primary"]

        # 그림자 효과
        if show_shadow:
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.04), Inches(y + 0.04),
                Inches(width), Inches(height)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(210, 210, 210)
            shadow.line.fill.background()

        # 메인 카드
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.design.BRAND_COLORS["white"]
        card.line.color.rgb = RGBColor(230, 230, 230)
        card.line.width = Pt(1)

        # 좌측 컬러 바
        color_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(0.06), Inches(height)
        )
        color_bar.fill.solid()
        color_bar.fill.fore_color.rgb = accent
        color_bar.line.fill.background()

        # 상단 원형 아이콘
        icon_size = 0.7
        icon_x = x + (width - icon_size) / 2
        icon_y_pos = y + 0.15

        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(icon_x), Inches(icon_y_pos),
            Inches(icon_size), Inches(icon_size)
        )
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = accent
        icon_bg.line.fill.background()

        icon_text = self._get_icon_text(icon, card_index)
        tf = icon_bg.text_frame
        tf.word_wrap = False
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = icon_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["white"]

        # 제목
        title_y = icon_y_pos + icon_size + 0.1
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(title_y),
            Inches(width - 0.3), Inches(0.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = accent

        # 구분선
        divider_y = title_y + 0.5
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.2), Inches(divider_y),
            Inches(width - 0.4), Inches(0.015)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(230, 230, 230)
        divider.line.fill.background()

        # 내용
        content_y = divider_y + 0.1
        content_h = height - (content_y - y) - 0.15
        content_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(content_y),
            Inches(width - 0.4), Inches(content_h)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = content
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.line_spacing = 1.2

        return card

    def _add_card_gradient(
        self, slide, title: str, content: str, x: float, y: float,
        width: float, height: float, accent_color: RGBColor = None,
        show_shadow: bool = True, icon: str = None, card_index: int = 0
    ):
        """[개선] 그라데이션 헤더 카드 - ref_1.pptx 스타일 기반"""
        accent = accent_color or self.design.BRAND_COLORS["primary"]
        header_height = 0.75  # 헤더 높이 조정

        # 부드러운 그림자 (더 넓고 연하게)
        if show_shadow:
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.04), Inches(y + 0.04),
                Inches(width), Inches(height)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(180, 190, 200)
            shadow.line.fill.background()

        # 메인 카드 배경 (흰색)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.design.BRAND_COLORS["white"]
        card.line.color.rgb = self.design.BRAND_COLORS["card_border"]
        card.line.width = Pt(1)

        # 상단 컬러 헤더 영역
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(header_height)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = accent
        header.line.fill.background()

        # 헤더 내 아이콘 (둥근 사각형)
        icon_size = 0.45
        icon_x = x + 0.12
        icon_y_pos = y + (header_height - icon_size) / 2

        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(icon_x), Inches(icon_y_pos),
            Inches(icon_size), Inches(icon_size)
        )
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = self.design.BRAND_COLORS["white"]
        icon_bg.line.fill.background()

        icon_text = self._get_icon_text(icon, card_index)
        tf = icon_bg.text_frame
        tf.word_wrap = False
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = icon_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = accent

        # 헤더 내 제목 (흰색)
        title_x = icon_x + icon_size + 0.1
        title_box = slide.shapes.add_textbox(
            Inches(title_x), Inches(y + 0.08),
            Inches(width - (title_x - x) - 0.1), Inches(header_height - 0.16)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT
        p.font.name = self.design.FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["white"]

        # 본문 영역 배경 (연한 파랑 - ref_1.pptx 스타일)
        content_y = y + header_height
        content_h = height - header_height
        content_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.05), Inches(content_y + 0.05),
            Inches(width - 0.1), Inches(content_h - 0.1)
        )
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = self.design.BRAND_COLORS["content_box"]
        content_bg.line.fill.background()

        # 본문 내용
        content_box = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(content_y + 0.12),
            Inches(width - 0.3), Inches(content_h - 0.2)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = content
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = self.design.BRAND_COLORS["text"]
        p.line_spacing = 1.4

        return card

    def _add_card_modern(
        self, slide, title: str, content: str, x: float, y: float,
        width: float, height: float, accent_color: RGBColor = None,
        show_shadow: bool = True, icon: str = None, card_index: int = 0
    ):
        """[개선] 모던 카드 - ref_1.pptx 기반 좌측 아이콘 강조형"""
        accent = accent_color or self.design.BRAND_COLORS["primary"]
        icon_area_width = 0.9

        # 부드러운 그림자
        if show_shadow:
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.04), Inches(y + 0.04),
                Inches(width), Inches(height)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(180, 190, 200)
            shadow.line.fill.background()

        # 메인 카드 (콘텐츠 박스 색상)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.design.BRAND_COLORS["content_box"]
        card.line.color.rgb = RGBColor(200, 210, 225)
        card.line.width = Pt(1)

        # 좌측 아이콘 영역 배경 (악센트 색상)
        icon_bg_rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(icon_area_width), Inches(height)
        )
        icon_bg_rect.fill.solid()
        icon_bg_rect.fill.fore_color.rgb = accent
        icon_bg_rect.line.fill.background()

        # 큰 아이콘 (원형, 흰색 배경)
        icon_size = 0.55
        icon_x = x + (icon_area_width - icon_size) / 2
        icon_y_pos = y + (height - icon_size) / 2

        icon_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(icon_x), Inches(icon_y_pos),
            Inches(icon_size), Inches(icon_size)
        )
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = self.design.BRAND_COLORS["white"]
        icon_circle.line.fill.background()

        icon_text = self._get_icon_text(icon, card_index)
        tf = icon_circle.text_frame
        tf.word_wrap = False
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = icon_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = accent

        # 우측 콘텐츠 영역
        content_x = x + icon_area_width + 0.12
        content_width = width - icon_area_width - 0.2

        # 제목
        title_box = slide.shapes.add_textbox(
            Inches(content_x), Inches(y + 0.15),
            Inches(content_width), Inches(0.45)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = self.design.FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["primary"]

        # 구분선
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(content_x), Inches(y + 0.6),
            Inches(content_width - 0.1), Inches(0.02)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self.design.BRAND_COLORS["secondary"]
        divider.line.fill.background()

        # 내용
        content_box = slide.shapes.add_textbox(
            Inches(content_x), Inches(y + 0.7),
            Inches(content_width), Inches(height - 0.85)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = content
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = self.design.BRAND_COLORS["text"]
        p.line_spacing = 1.3

        return card

    def _add_card_solid(
        self, slide, title: str, content: str, x: float, y: float,
        width: float, height: float, accent_color: RGBColor = None,
        show_shadow: bool = True, icon: str = None, card_index: int = 0
    ):
        """[개선] 솔리드 카드 - ref_1.pptx 기반 전체 컬러 배경"""
        accent = accent_color or self.design.BRAND_COLORS["secondary"]  # #4F81BD 사용
        highlight = self.design.BRAND_COLORS.get("highlight", RGBColor(255, 192, 0))

        # 부드러운 그림자
        if show_shadow:
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.04), Inches(y + 0.04),
                Inches(width), Inches(height)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(150, 160, 175)
            shadow.line.fill.background()

        # 메인 카드 (전체 컬러)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = accent
        card.line.fill.background()

        # 상단 아이콘 영역 (흰색 원형)
        icon_size = 0.5
        icon_x = x + (width - icon_size) / 2
        icon_y_pos = y + 0.18

        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(icon_x), Inches(icon_y_pos),
            Inches(icon_size), Inches(icon_size)
        )
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = self.design.BRAND_COLORS["white"]
        icon_bg.line.fill.background()

        icon_text = self._get_icon_text(icon, card_index)
        tf = icon_bg.text_frame
        tf.word_wrap = False
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = icon_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = accent

        # 제목 (흰색, 굵게)
        title_y = icon_y_pos + icon_size + 0.1
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(title_y),
            Inches(width - 0.2), Inches(0.45)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["white"]

        # 내용 (흰색, 약간 투명한 배경)
        content_y = title_y + 0.5
        content_h = height - (content_y - y) - 0.12

        # 내용 배경 박스 (반투명 효과)
        content_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + 0.08), Inches(content_y),
            Inches(width - 0.16), Inches(content_h)
        )
        # 약간 어두운 악센트 색상
        darker = RGBColor(
            max(0, accent[0] - 25),
            max(0, accent[1] - 25),
            max(0, accent[2] - 25)
        )
        content_bg.fill.solid()
        content_bg.fill.fore_color.rgb = darker
        content_bg.line.fill.background()

        content_box = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(content_y + 0.08),
            Inches(width - 0.3), Inches(content_h - 0.12)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = content
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = self.design.BRAND_COLORS["white"]
        p.line_spacing = 1.3

        return card

    def _add_card_outline(
        self, slide, title: str, content: str, x: float, y: float,
        width: float, height: float, accent_color: RGBColor = None,
        show_shadow: bool = True, icon: str = None, card_index: int = 0
    ):
        """[개선] 아웃라인 카드 - ref_1.pptx 기반 테두리 강조"""
        accent = accent_color or self.design.BRAND_COLORS["secondary"]  # 미드블루

        # 부드러운 그림자
        if show_shadow:
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.03), Inches(y + 0.03),
                Inches(width), Inches(height)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(180, 190, 200)
            shadow.line.fill.background()

        # 메인 카드 (두꺼운 컬러 테두리)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.design.BRAND_COLORS["content_box"]
        card.line.color.rgb = accent
        card.line.width = Pt(2.5)

        # 상단 아이콘 (테두리 원형)
        icon_size = 0.5
        icon_x = x + (width - icon_size) / 2
        icon_y_pos = y + 0.15

        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(icon_x), Inches(icon_y_pos),
            Inches(icon_size), Inches(icon_size)
        )
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = self.design.BRAND_COLORS["white"]
        icon_bg.line.color.rgb = accent
        icon_bg.line.width = Pt(2)

        icon_text = self._get_icon_text(icon, card_index)
        tf = icon_bg.text_frame
        tf.word_wrap = False
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = icon_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = accent

        # 제목
        title_y = icon_y_pos + icon_size + 0.08
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.12), Inches(title_y),
            Inches(width - 0.24), Inches(0.4)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["primary"]

        # 내용
        content_y = title_y + 0.45
        content_h = height - (content_y - y) - 0.12
        content_box = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(content_y),
            Inches(width - 0.3), Inches(content_h)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = content
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = self.design.BRAND_COLORS["text"]
        p.line_spacing = 1.3

        return card

    def _add_card_minimal(
        self, slide, title: str, content: str, x: float, y: float,
        width: float, height: float, accent_color: RGBColor = None,
        show_shadow: bool = True, icon: str = None, card_index: int = 0
    ):
        """[개선] 미니멀 카드 - ref_1.pptx 기반 하단 컬러 라인"""
        accent = accent_color or self.design.BRAND_COLORS["accent"]  # 딥블루

        # 가벼운 그림자
        if show_shadow:
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + 0.02), Inches(y + 0.02),
                Inches(width), Inches(height)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(200, 210, 220)
            shadow.line.fill.background()

        # 메인 카드 (콘텐츠 박스 배경)
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.design.BRAND_COLORS["content_box"]
        card.line.color.rgb = RGBColor(200, 210, 225)
        card.line.width = Pt(1)

        # 하단 컬러 라인
        bottom_line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y + height - 0.06),
            Inches(width), Inches(0.06)
        )
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = accent
        bottom_line.line.fill.background()

        # 좌측 상단 작은 아이콘 (사각형)
        icon_size = 0.4
        icon_x = x + 0.12
        icon_y_pos = y + 0.12

        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(icon_x), Inches(icon_y_pos),
            Inches(icon_size), Inches(icon_size)
        )
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = accent
        icon_bg.line.fill.background()

        icon_text = self._get_icon_text(icon, card_index)
        tf = icon_bg.text_frame
        tf.word_wrap = False
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = icon_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["white"]

        # 제목 (아이콘 옆)
        title_box = slide.shapes.add_textbox(
            Inches(icon_x + icon_size + 0.08), Inches(icon_y_pos),
            Inches(width - icon_size - 0.35), Inches(icon_size)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT
        p.font.name = self.design.FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["primary"]

        # 내용
        content_y = icon_y_pos + icon_size + 0.1
        content_h = height - (content_y - y) - 0.15
        content_box = slide.shapes.add_textbox(
            Inches(x + 0.12), Inches(content_y),
            Inches(width - 0.24), Inches(content_h)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = content
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = self.design.BRAND_COLORS["text"]
        p.line_spacing = 1.3

        return card

    def _add_card_banner(
        self, slide, title: str, content: str, x: float, y: float,
        width: float, height: float, accent_color: RGBColor = None,
        show_shadow: bool = True, icon: str = None, card_index: int = 0
    ):
        """[개선] 배너 카드 - ref_1.pptx 기반 상단 풀 컬러 배너"""
        accent = accent_color or self.design.BRAND_COLORS["success"]  # 티일 그린
        banner_height = 0.6

        # 그림자
        if show_shadow:
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.04), Inches(y + 0.04),
                Inches(width), Inches(height)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(180, 190, 200)
            shadow.line.fill.background()

        # 메인 카드 (콘텐츠 박스 배경)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.design.BRAND_COLORS["content_box"]
        card.line.color.rgb = RGBColor(200, 210, 225)
        card.line.width = Pt(1)

        # 상단 배너
        banner = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(banner_height)
        )
        banner.fill.solid()
        banner.fill.fore_color.rgb = accent
        banner.line.fill.background()

        # 배너 위 큰 아이콘 (중앙, 배너 아래로 튀어나옴)
        icon_size = 0.65
        icon_x = x + (width - icon_size) / 2
        icon_y_pos = y + banner_height - icon_size / 2

        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(icon_x), Inches(icon_y_pos),
            Inches(icon_size), Inches(icon_size)
        )
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = self.design.BRAND_COLORS["white"]
        icon_bg.line.color.rgb = accent
        icon_bg.line.width = Pt(2)

        icon_text = self._get_icon_text(icon, card_index)
        tf = icon_bg.text_frame
        tf.word_wrap = False
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = icon_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = accent

        # 제목
        title_y = icon_y_pos + icon_size + 0.08
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(title_y),
            Inches(width - 0.2), Inches(0.4)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["primary"]

        # 내용
        content_y = title_y + 0.45
        content_h = height - (content_y - y) - 0.1
        content_box = slide.shapes.add_textbox(
            Inches(x + 0.12), Inches(content_y),
            Inches(width - 0.24), Inches(content_h)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = content
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = self.design.BRAND_COLORS["text"]
        p.line_spacing = 1.3

        return card

    def _add_card_split(
        self, slide, title: str, content: str, x: float, y: float,
        width: float, height: float, accent_color: RGBColor = None,
        show_shadow: bool = True, icon: str = None, card_index: int = 0
    ):
        """[개선] 스플릿 카드 - ref_1.pptx 기반 상하 분할"""
        accent = accent_color or self.design.BRAND_COLORS["teal"]  # 다크 티일
        split_ratio = 0.38  # 상단 38%

        # 그림자
        if show_shadow:
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.03), Inches(y + 0.03),
                Inches(width), Inches(height)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(180, 190, 200)
            shadow.line.fill.background()

        # 하단 영역 (콘텐츠 박스 색상)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.design.BRAND_COLORS["content_box"]
        card.line.color.rgb = RGBColor(200, 210, 225)
        card.line.width = Pt(1)

        # 상단 컬러 영역
        top_height = height * split_ratio
        top_area = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(top_height)
        )
        top_area.fill.solid()
        top_area.fill.fore_color.rgb = accent
        top_area.line.fill.background()

        # 아이콘 (상단 영역 중앙, 원형)
        icon_size = 0.5
        icon_x = x + (width - icon_size) / 2
        icon_y_pos = y + (top_height - icon_size) / 2

        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(icon_x), Inches(icon_y_pos),
            Inches(icon_size), Inches(icon_size)
        )
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = self.design.BRAND_COLORS["white"]
        icon_bg.line.fill.background()

        icon_text = self._get_icon_text(icon, card_index)
        tf = icon_bg.text_frame
        tf.word_wrap = False
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = icon_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = accent

        # 제목 (분할선 바로 아래)
        title_y = y + top_height + 0.08
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(title_y),
            Inches(width - 0.2), Inches(0.4)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["primary"]

        # 내용
        content_y = title_y + 0.45
        content_h = height - (content_y - y) - 0.08
        content_box = slide.shapes.add_textbox(
            Inches(x + 0.12), Inches(content_y),
            Inches(width - 0.24), Inches(content_h)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = content
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = self.design.BRAND_COLORS["text"]
        p.line_spacing = 1.3

        return card

    def _add_card_accent(
        self, slide, title: str, content: str, x: float, y: float,
        width: float, height: float, accent_color: RGBColor = None,
        show_shadow: bool = True, icon: str = None, card_index: int = 0
    ):
        """[개선] 악센트 카드 - ref_1.pptx 기반 좌측 악센트 바"""
        accent = accent_color or self.design.BRAND_COLORS["danger"]  # 소프트 레드
        accent_bar_width = 0.1

        # 그림자
        if show_shadow:
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.03), Inches(y + 0.03),
                Inches(width), Inches(height)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(180, 190, 200)
            shadow.line.fill.background()

        # 메인 카드 (콘텐츠 박스 배경)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.design.BRAND_COLORS["content_box"]
        card.line.color.rgb = RGBColor(200, 210, 225)
        card.line.width = Pt(1)

        # 좌측 두꺼운 악센트 바
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(accent_bar_width), Inches(height)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()

        # 큰 아이콘 (악센트 바 옆, 상단)
        icon_size = 0.55
        icon_x = x + accent_bar_width + 0.1
        icon_y_pos = y + 0.12

        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(icon_x), Inches(icon_y_pos),
            Inches(icon_size), Inches(icon_size)
        )
        icon_bg.fill.solid()
        icon_bg.fill.fore_color.rgb = accent
        icon_bg.line.fill.background()

        icon_text = self._get_icon_text(icon, card_index)
        tf = icon_bg.text_frame
        tf.word_wrap = False
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = icon_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["white"]

        # 제목 (아이콘 옆)
        title_box = slide.shapes.add_textbox(
            Inches(icon_x + icon_size + 0.08), Inches(icon_y_pos),
            Inches(width - icon_x - icon_size - 0.15 + x), Inches(icon_size)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.LEFT
        p.font.name = self.design.FONT_TITLE
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["primary"]

        # 내용
        content_y = icon_y_pos + icon_size + 0.08
        content_h = height - (content_y - y) - 0.1
        content_box = slide.shapes.add_textbox(
            Inches(x + accent_bar_width + 0.12), Inches(content_y),
            Inches(width - accent_bar_width - 0.22), Inches(content_h)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = content
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = self.design.BRAND_COLORS["text"]
        p.line_spacing = 1.3

        return card

    # 고정 작성자 정보
    FIXED_AUTHOR = "미래융합설계센터 알고리즘개발팀 강민규 선임"

    def add_cover_slide(
        self,
        title: str,
        date: Optional[str] = None,
        author: str = "",
        report_type: str = "정보공유"
    ):
        """표지 슬라이드 추가 (현대하모니 폰트 적용)"""
        layout = self._get_layout("제목 슬라이드")
        slide = self.prs.slides.add_slide(layout)

        # 제목 설정 - 현대하모니M 48pt 굵게 (흰색)
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(48)
            para.font.bold = True
            # 색상 지정하지 않음 - 템플릿 기본값 사용 (흰색)

        # 날짜 설정 - 현대하모니L 24pt
        if date is None:
            date = datetime.now().strftime("%Y. %m. %d")

        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = date
                for para in shape.text_frame.paragraphs:
                    para.font.name = self.design.FONT_BODY
                    para.font.size = Pt(24)
                    para.font.bold = False
                    # 색상 지정하지 않음 - 템플릿 기본값 사용
                break

        # 작성자 정보 - 현대하모니L 24pt 굵게 (흰색)
        txBox = slide.shapes.add_textbox(
            Inches(1.76), Inches(6.55),
            Inches(7.31), Inches(0.48)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = self.FIXED_AUTHOR
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(24)
        p.font.bold = True
        # 색상 지정하지 않음 - 템플릿 기본값 사용 (흰색)
        p.alignment = PP_ALIGN.CENTER

        # 보고 유형 - 현대하모니L 14pt (박스 없이 텍스트만)
        report_type_text = self.REPORT_TYPES.get(report_type, self.REPORT_TYPES["정보공유"])
        txBox = slide.shapes.add_textbox(
            Inches(8.5), Inches(0.35),
            Inches(1.8), Inches(0.35)
        )
        tf = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = report_type_text
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = False
        p.font.color.rgb = RGBColor(255, 255, 255)  # 흰색 텍스트
        p.alignment = PP_ALIGN.RIGHT

        # 표지는 제목 플레이스홀더 사용
        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def add_section_slide(
        self,
        section_number: str,
        section_title: str,
        subtitle: str = ""
    ):
        """섹션 구분 슬라이드 추가"""
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 기존 제목 숨기기
        if slide.shapes.title:
            slide.shapes.title.text = ""

        # 배경 악센트 (좌측 세로 바)
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.25), Inches(7.5)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = self.design.BRAND_COLORS["primary"]
        accent_bar.line.fill.background()

        # 섹션 번호 (큰 원)
        self._add_icon_box(
            slide, section_number,
            x=1.0, y=2.5, size=1.2,
            bg_color=self.design.BRAND_COLORS["primary"]
        )

        # 섹션 제목
        title_box = slide.shapes.add_textbox(
            Inches(2.5), Inches(2.6),
            Inches(7.0), Inches(1.0)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["primary"]

        # 부제목
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(2.5), Inches(3.7),
                Inches(7.0), Inches(0.5)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(64, 64, 64)

        self._clear_unused_placeholders(slide)
        return slide

    def add_content_slide(
        self,
        title: str,
        content: List[str],
        layout_name: str = "제목 및 내용"
    ):
        """내용 슬라이드 추가"""
        layout = self._get_layout(layout_name)
        slide = self.prs.slides.add_slide(layout)

        # 제목 설정 (진한 검정)
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 본문 내용
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 10:
                tf = shape.text_frame
                tf.clear()

                for i, item in enumerate(content):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    p.text = item
                    p.font.name = self.design.FONT_BODY
                    p.font.size = Pt(16)
                    p.font.color.rgb = self.design.BRAND_COLORS["black"]
                    p.level = 0
                    p.space_after = Pt(12)
                break

        # 제목(0)과 본문(10) 플레이스홀더는 사용 중
        self._clear_unused_placeholders(slide, used_placeholder_idx=[0, 10])
        return slide

    def add_content_boxed_slide(
        self,
        title: str,
        sections: List[Dict[str, Any]],
        columns: int = 1
    ):
        """소주제별 박스로 구분된 콘텐츠 슬라이드 - 각 항목별 개별 박스

        Args:
            title: 슬라이드 제목
            sections: 소주제 목록
                - title: 소주제 제목
                - items: 내용 리스트
                - color: 악센트 색상 (선택, primary/secondary/accent/success/warning/danger)
            columns: 열 개수 (1 또는 2)

        Example YAML:
            - type: content_boxed
              title: "주요 기능"
              columns: 2
              sections:
                - title: "데이터 처리"
                  items:
                    - "실시간 스트리밍 처리"
                    - "배치 처리 지원"
                  color: "primary"
                - title: "보안 기능"
                  items:
                    - "암호화 통신"
                    - "접근 제어"
                  color: "secondary"
        """
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 제목 설정
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 레이아웃 계산
        margin_left = 0.4
        start_y = 1.1
        total_width = 10.0 - margin_left * 2
        section_spacing = 0.15
        item_spacing = 0.08

        num_sections = len(sections)
        if num_sections == 0:
            self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
            return slide

        # 색상 팔레트 - 테두리용 연한 색상 포함
        color_map = {
            "primary": (self.design.BRAND_COLORS["primary"], RGBColor(200, 215, 235)),
            "secondary": (self.design.BRAND_COLORS["secondary"], RGBColor(200, 225, 245)),
            "accent": (self.design.BRAND_COLORS["accent"], RGBColor(210, 220, 240)),
            "success": (self.design.BRAND_COLORS["success"], RGBColor(200, 235, 200)),
            "warning": (self.design.BRAND_COLORS["warning"], RGBColor(255, 235, 200)),
            "danger": (self.design.BRAND_COLORS["danger"], RGBColor(245, 210, 210)),
            "highlight": (self.design.BRAND_COLORS["highlight"], RGBColor(255, 240, 200)),
        }

        # 열 배치 계산
        if columns == 2 and num_sections >= 2:
            col_width = (total_width - section_spacing) / 2
        else:
            columns = 1
            col_width = total_width

        # 각 열의 현재 Y 위치 추적
        col_y = [start_y] * columns

        # 각 섹션 박스 생성
        for idx, section in enumerate(sections):
            col = idx % columns if columns == 2 else 0
            x = margin_left + col * (col_width + section_spacing)
            y = col_y[col]

            # 악센트 색상
            color_name = section.get("color", "primary")
            accent_color, border_color = color_map.get(color_name, color_map["primary"])

            # 소주제 제목
            section_title = section.get("title", "")
            items = section.get("items", [])

            # 제목 높이
            title_height = 0.35
            item_box_height = 0.32

            # 소주제 제목 박스 (그림자 효과)
            self._add_shadow_box(
                slide,
                x=x, y=y,
                width=col_width, height=title_height,
                fill_color=accent_color,
                border_color=accent_color,
                shadow_offset=0.03,
                shadow_color=RGBColor(180, 180, 180)
            )

            # 소주제 제목 텍스트
            title_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(y + 0.05),
                Inches(col_width - 0.3), Inches(title_height - 0.1)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            tf.anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = section_title
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.design.BRAND_COLORS["white"]

            y += title_height + item_spacing

            # 각 항목별 개별 박스 생성
            for i, item in enumerate(items):
                # 항목 박스 (그림자 효과 + 테두리)
                item_box = self._add_shadow_box(
                    slide,
                    x=x, y=y,
                    width=col_width, height=item_box_height,
                    fill_color=self.design.BRAND_COLORS["white"],
                    border_color=border_color,
                    shadow_offset=0.025,
                    shadow_color=RGBColor(220, 220, 220)
                )

                # 좌측 악센트 바 (작은 크기)
                accent_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x + 0.02), Inches(y + 0.06),
                    Inches(0.04), Inches(item_box_height - 0.12)
                )
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = accent_color
                accent_bar.line.fill.background()

                # 항목 텍스트
                text_box = slide.shapes.add_textbox(
                    Inches(x + 0.12), Inches(y + 0.04),
                    Inches(col_width - 0.2), Inches(item_box_height - 0.08)
                )
                tf = text_box.text_frame
                tf.word_wrap = True
                tf.anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = item
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(11)
                p.font.color.rgb = self.design.BRAND_COLORS["text"]

                y += item_box_height + item_spacing

            # 섹션 간 여백 추가
            y += section_spacing
            col_y[col] = y

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def add_content_with_icons(
        self,
        title: str,
        items: List[Dict[str, str]],
    ):
        """아이콘이 있는 내용 슬라이드 (번호/기호 + 텍스트)"""
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 제목 (진한 검정)
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 아이템들
        start_y = 1.0
        for i, item in enumerate(items):
            y = start_y + i * 0.9

            # 아이콘/번호
            icon = item.get("icon", str(i + 1))
            color_idx = i % len(self.design.GRADIENT_BLUE)
            self._add_icon_box(
                slide, icon,
                x=0.5, y=y,
                size=0.5,
                bg_color=self.design.GRADIENT_BLUE[color_idx]
            )

            # 제목 (검정)
            title_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(y),
                Inches(8.5), Inches(0.4)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = item.get("title", "")
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.design.BRAND_COLORS["black"]

            # 설명 (진한 회색)
            if item.get("description"):
                desc_box = slide.shapes.add_textbox(
                    Inches(1.2), Inches(y + 0.4),
                    Inches(8.5), Inches(0.4)
                )
                tf = desc_box.text_frame
                p = tf.paragraphs[0]
                p.text = item["description"]
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(64, 64, 64)

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def add_comparison_slide(
        self,
        title: str,
        left_title: str,
        left_items: List[str],
        right_title: str,
        right_items: List[str],
        left_color: str = "primary",
        right_color: str = "secondary"
    ):
        """비교 슬라이드 (좌우 대비) - 큰 바운더리 박스 + 텍스트 제목"""
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 슬라이드 제목 (진한 검정)
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 색상 가져오기
        left_color_rgb = self.design.BRAND_COLORS.get(left_color, self.design.BRAND_COLORS["primary"])
        right_color_rgb = self.design.BRAND_COLORS.get(right_color, self.design.BRAND_COLORS["secondary"])

        # 레이아웃 설정
        box_y = 1.1
        box_height = 5.6
        left_x = 0.4
        right_x = 5.5
        box_width = 4.9
        title_height = 0.5
        content_start_y = box_y + title_height + 0.15

        # ===== 좌측 영역 =====
        # 좌측 바운더리 박스 (큰 테두리)
        left_boundary = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left_x), Inches(box_y),
            Inches(box_width), Inches(box_height)
        )
        left_boundary.fill.solid()
        left_boundary.fill.fore_color.rgb = RGBColor(248, 250, 255)
        left_boundary.line.color.rgb = left_color_rgb
        left_boundary.line.width = Pt(2)

        # 좌측 제목 배경 (바운더리 내부 상단)
        left_title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left_x), Inches(box_y),
            Inches(box_width), Inches(title_height)
        )
        left_title_bg.fill.solid()
        left_title_bg.fill.fore_color.rgb = left_color_rgb
        left_title_bg.line.fill.background()

        # 좌측 제목 텍스트
        left_title_box = slide.shapes.add_textbox(
            Inches(left_x + 0.1), Inches(box_y + 0.05),
            Inches(box_width - 0.2), Inches(title_height - 0.1)
        )
        tf = left_title_box.text_frame
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = left_title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["white"]

        # 좌측 아이템들 - 바운더리 내부에 표시
        item_height = 0.75
        item_gap = 0.08
        max_items = min(len(left_items), 6)

        for i, item in enumerate(left_items[:max_items]):
            item_y = content_start_y + i * (item_height + item_gap)

            # 아이템 박스
            item_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left_x + 0.15), Inches(item_y),
                Inches(box_width - 0.3), Inches(item_height)
            )
            item_bg.fill.solid()
            item_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
            item_bg.line.color.rgb = RGBColor(220, 225, 235)
            item_bg.line.width = Pt(1)

            # 좌측 컬러바
            color_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left_x + 0.15), Inches(item_y),
                Inches(0.05), Inches(item_height)
            )
            color_bar.fill.solid()
            color_bar.fill.fore_color.rgb = left_color_rgb
            color_bar.line.fill.background()

            # 아이템 텍스트
            item_box = slide.shapes.add_textbox(
                Inches(left_x + 0.3), Inches(item_y + 0.08),
                Inches(box_width - 0.5), Inches(item_height - 0.16)
            )
            tf = item_box.text_frame
            tf.word_wrap = True
            tf.anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(12)
            p.font.color.rgb = self.design.BRAND_COLORS["dark"]

        # ===== 우측 영역 =====
        # 우측 바운더리 박스 (큰 테두리)
        right_boundary = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(right_x), Inches(box_y),
            Inches(box_width), Inches(box_height)
        )
        right_boundary.fill.solid()
        right_boundary.fill.fore_color.rgb = RGBColor(248, 255, 250)
        right_boundary.line.color.rgb = right_color_rgb
        right_boundary.line.width = Pt(2)

        # 우측 제목 배경 (바운더리 내부 상단)
        right_title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(right_x), Inches(box_y),
            Inches(box_width), Inches(title_height)
        )
        right_title_bg.fill.solid()
        right_title_bg.fill.fore_color.rgb = right_color_rgb
        right_title_bg.line.fill.background()

        # 우측 제목 텍스트
        right_title_box = slide.shapes.add_textbox(
            Inches(right_x + 0.1), Inches(box_y + 0.05),
            Inches(box_width - 0.2), Inches(title_height - 0.1)
        )
        tf = right_title_box.text_frame
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = right_title
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["white"]

        # 우측 아이템들 - 바운더리 내부에 표시
        max_items = min(len(right_items), 6)

        for i, item in enumerate(right_items[:max_items]):
            item_y = content_start_y + i * (item_height + item_gap)

            # 아이템 박스
            item_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(right_x + 0.15), Inches(item_y),
                Inches(box_width - 0.3), Inches(item_height)
            )
            item_bg.fill.solid()
            item_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
            item_bg.line.color.rgb = RGBColor(220, 235, 225)
            item_bg.line.width = Pt(1)

            # 좌측 컬러바
            color_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(right_x + 0.15), Inches(item_y),
                Inches(0.05), Inches(item_height)
            )
            color_bar.fill.solid()
            color_bar.fill.fore_color.rgb = right_color_rgb
            color_bar.line.fill.background()

            # 아이템 텍스트
            item_box = slide.shapes.add_textbox(
                Inches(right_x + 0.3), Inches(item_y + 0.08),
                Inches(box_width - 0.5), Inches(item_height - 0.16)
            )
            tf = item_box.text_frame
            tf.word_wrap = True
            tf.anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(12)
            p.font.color.rgb = self.design.BRAND_COLORS["dark"]

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide
        tf.word_wrap = False
        tf.anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = "VS"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["white"]

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def add_text_slide(
        self,
        title: str,
        text_blocks: List[Dict[str, Any]]
    ):
        """자유 형식 텍스트 슬라이드"""
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 제목 (진한 검정)
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 텍스트 블록
        for block in text_blocks:
            left = Inches(block.get("x", 0.5))
            top = Inches(block.get("y", 1.0))
            width = Inches(block.get("width", 9.0))
            height = Inches(block.get("height", 1.0))

            # 배경 박스 (선택적)
            if block.get("background"):
                bg_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left, top, width, height
                )
                bg_color = block.get("bg_color", "light")
                if isinstance(bg_color, str) and bg_color in self.design.BRAND_COLORS:
                    bg_box.fill.solid()
                    bg_box.fill.fore_color.rgb = self.design.BRAND_COLORS[bg_color]
                bg_box.line.fill.background()

            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = block.get("text", "")
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(block.get("font_size", 14))
            p.font.bold = block.get("bold", False)

            # 텍스트 색상 (기본 검정)
            text_color = block.get("color", "black")
            if isinstance(text_color, str) and text_color in self.design.BRAND_COLORS:
                p.font.color.rgb = self.design.BRAND_COLORS[text_color]
            else:
                p.font.color.rgb = self.design.BRAND_COLORS["black"]

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def add_table_slide(
        self,
        title: str,
        headers: List[str],
        rows: List[List[str]],
        col_widths: Optional[List[float]] = None,
        highlight_rows: List[int] = None
    ):
        """표 슬라이드"""
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 제목 (진한 검정)
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 표 생성
        num_rows = len(rows) + 1
        num_cols = len(headers)

        left = Inches(0.4)
        top = Inches(1.2)
        width = Inches(10.0)
        height = Inches(0.45 * num_rows)

        table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

        # 열 너비
        if col_widths:
            for i, w in enumerate(col_widths):
                table.columns[i].width = Inches(w)
        else:
            col_width = 10.0 / num_cols
            for col in table.columns:
                col.width = Inches(col_width)

        # 헤더 스타일
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.design.BRAND_COLORS["primary"]

            for para in cell.text_frame.paragraphs:
                para.font.name = self.design.FONT_BODY
                para.font.size = Pt(12)
                para.font.bold = True
                para.font.color.rgb = self.design.BRAND_COLORS["white"]
                para.alignment = PP_ALIGN.CENTER

            cell.text_frame.anchor = MSO_ANCHOR.MIDDLE

        # 데이터 행
        highlight_rows = highlight_rows or []
        for row_idx, row_data in enumerate(rows):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_text)

                # 줄무늬 배경
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(248, 248, 248)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.design.BRAND_COLORS["white"]

                # 강조 행
                if row_idx in highlight_rows:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 200)

                for para in cell.text_frame.paragraphs:
                    para.font.name = self.design.FONT_BODY
                    para.font.size = Pt(11)
                    para.font.color.rgb = self.design.BRAND_COLORS["black"]
                    para.alignment = PP_ALIGN.CENTER

                cell.text_frame.anchor = MSO_ANCHOR.MIDDLE

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def add_cards_slide(
        self,
        title: str,
        cards: List[Dict[str, str]],
        columns: int = 3,
        card_style: str = None
    ):
        """카드 그리드 슬라이드 - 프로페셔널 디자인

        Args:
            title: 슬라이드 제목
            cards: 카드 목록 [{title, content, icon}]
            columns: 열 개수 (1-4)
            card_style: 카드 스타일 (classic, gradient, modern, solid)
                       None이면 전역 설정 사용
        """
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 제목 (진한 검정)
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 카드 배치 - 슬라이드 영역 내로 제한
        num_cards = len(cards)
        max_columns = min(columns, 4)  # 최대 4열
        actual_columns = min(max_columns, num_cards)
        rows = (num_cards + actual_columns - 1) // actual_columns

        # 슬라이드 영역 계산 (마진 고려)
        content_width = 9.6  # 좌우 마진 0.4씩
        content_height = 5.8  # 제목 영역 제외
        start_x = 0.4
        start_y = 1.2

        # 카드 크기 계산 (간격 포함)
        gap_x = 0.25
        gap_y = 0.25
        card_width = (content_width - (actual_columns - 1) * gap_x) / actual_columns
        card_height = min(2.2, (content_height - (rows - 1) * gap_y) / rows)

        for i, card in enumerate(cards):
            col = i % actual_columns
            row = i // actual_columns
            x = start_x + col * (card_width + gap_x)
            y = start_y + row * (card_height + gap_y)

            # 색상 순환
            color_keys = ["primary", "secondary", "accent", "success", "warning"]
            accent_color = self.design.BRAND_COLORS[color_keys[i % len(color_keys)]]

            self._add_card(
                slide,
                title=card.get("title", ""),
                content=card.get("content", ""),
                x=x, y=y,
                width=card_width,
                height=card_height,
                accent_color=accent_color,
                icon=card.get("icon"),
                card_index=i,
                card_style=card_style
            )

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def _add_box(
        self,
        slide,
        text: str,
        x: float,
        y: float,
        width: float,
        height: float,
        fill_color: RGBColor,
        text_color: RGBColor = None,
        font_size: int = 12,
        bold: bool = True,
        shape_type: MSO_SHAPE = MSO_SHAPE.ROUNDED_RECTANGLE,
        border_color: RGBColor = None,
        border_width: float = 1.0
    ):
        """박스 도형 추가"""
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )

        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = border_color or fill_color
        shape.line.width = Pt(border_width)

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = text_color or RGBColor(255, 255, 255)

        return shape

    def _add_arrow(
        self,
        slide,
        start_x: float,
        start_y: float,
        end_x: float,
        end_y: float,
        color: RGBColor = None,
        line_width: float = 2.0,
        dashed: bool = False
    ):
        """화살표 추가"""
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(start_x), Inches(start_y),
            Inches(end_x), Inches(end_y)
        )

        connector.line.color.rgb = color or RGBColor(89, 89, 89)
        connector.line.width = Pt(line_width)

        if dashed:
            connector.line.dash_style = 2

        connector.line._ln.append(
            connector.line._ln.makeelement(
                qn('a:tailEnd'),
                {'type': 'triangle', 'w': 'med', 'len': 'med'}
            )
        )

        return connector

    def _add_line(
        self,
        slide,
        start_x: float,
        start_y: float,
        end_x: float,
        end_y: float,
        color: RGBColor = None,
        line_width: float = 2.0
    ):
        """화살표 없는 선 추가"""
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(start_x), Inches(start_y),
            Inches(end_x), Inches(end_y)
        )

        connector.line.color.rgb = color or RGBColor(89, 89, 89)
        connector.line.width = Pt(line_width)

        return connector

    def _add_label(
        self,
        slide,
        text: str,
        x: float,
        y: float,
        font_size: int = 10,
        color: RGBColor = None,
        bold: bool = False
    ):
        """텍스트 레이블 추가"""
        txBox = slide.shapes.add_textbox(
            Inches(x), Inches(y),
            Inches(2), Inches(0.3)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color or RGBColor(89, 89, 89)

        return txBox

    def add_architecture_slide(
        self,
        title: str,
        components: List[Dict[str, Any]],
        connections: List[Dict[str, Any]] = None,
        labels: List[Dict[str, Any]] = None,
        palette: str = "blue"
    ):
        """아키텍처 다이어그램 슬라이드"""
        # 색상 팔레트
        COLOR_PALETTES = {
            "blue": {
                "primary": self.design.BRAND_COLORS["primary"],
                "secondary": self.design.BRAND_COLORS["secondary"],
                "accent": self.design.BRAND_COLORS["accent"],
                "light": self.design.BRAND_COLORS["light"],
                "text": self.design.BRAND_COLORS["white"],
                "dark_text": self.design.BRAND_COLORS["dark"],
                "green": self.design.BRAND_COLORS["success"],
                "orange": self.design.BRAND_COLORS["warning"],
                "purple": RGBColor(112, 48, 160),
                "gray": RGBColor(128, 128, 128),
            },
            "green": {
                "primary": RGBColor(0, 128, 64),
                "secondary": RGBColor(100, 180, 100),
                "accent": RGBColor(50, 100, 50),
                "light": RGBColor(230, 245, 230),
                "text": self.design.BRAND_COLORS["white"],
                "dark_text": self.design.BRAND_COLORS["dark"],
            },
            "purple": {
                "primary": RGBColor(102, 45, 145),
                "secondary": RGBColor(150, 100, 180),
                "accent": RGBColor(80, 35, 115),
                "light": RGBColor(240, 230, 250),
                "text": self.design.BRAND_COLORS["white"],
                "dark_text": self.design.BRAND_COLORS["dark"],
            },
        }

        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 제목 (진한 검정)
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        colors = COLOR_PALETTES.get(palette, COLOR_PALETTES["blue"])
        component_positions = {}

        # 도형 타입 매핑
        shape_types = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "diamond": MSO_SHAPE.DIAMOND,
            "parallelogram": MSO_SHAPE.PARALLELOGRAM,
            "hexagon": MSO_SHAPE.HEXAGON,
            "chevron": MSO_SHAPE.CHEVRON,
            "cylinder": MSO_SHAPE.CAN,
            "cloud": MSO_SHAPE.CLOUD,
            "document": MSO_SHAPE.FLOWCHART_DOCUMENT,
        }

        # 슬라이드 경계 정의 (widescreen 10" x 7.5", 제목 영역 제외)
        SLIDE_BOUNDS = {
            "min_x": 0.3,
            "max_x": 10.0,
            "min_y": 1.0,  # 제목 영역 아래
            "max_y": 7.0,
        }

        # 컴포넌트 경계 계산 (스케일링 필요 여부 확인)
        if components:
            all_x = [comp.get("x", 1.0) for comp in components]
            all_y = [comp.get("y", 1.0) for comp in components]
            all_right = [comp.get("x", 1.0) + comp.get("width", 2.0) for comp in components]
            all_bottom = [comp.get("y", 1.0) + comp.get("height", 0.8) for comp in components]

            content_min_x = min(all_x)
            content_max_x = max(all_right)
            content_min_y = min(all_y)
            content_max_y = max(all_bottom)

            content_width = content_max_x - content_min_x
            content_height = content_max_y - content_min_y

            available_width = SLIDE_BOUNDS["max_x"] - SLIDE_BOUNDS["min_x"]
            available_height = SLIDE_BOUNDS["max_y"] - SLIDE_BOUNDS["min_y"]

            # 스케일 계산 (경계를 벗어나면 축소)
            scale_x = min(1.0, available_width / content_width) if content_width > 0 else 1.0
            scale_y = min(1.0, available_height / content_height) if content_height > 0 else 1.0
            scale = min(scale_x, scale_y)

            # 오프셋 계산 (중앙 정렬)
            offset_x = SLIDE_BOUNDS["min_x"] - content_min_x * scale + (available_width - content_width * scale) / 2
            offset_y = SLIDE_BOUNDS["min_y"] - content_min_y * scale + (available_height - content_height * scale) / 2
        else:
            scale = 1.0
            offset_x = 0
            offset_y = 0

        # 컴포넌트 그리기 (스케일 및 오프셋 적용)
        for comp in components:
            comp_id = comp.get("id", comp.get("name", comp.get("text", "")))
            orig_x = comp.get("x", 1.0)
            orig_y = comp.get("y", 1.0)
            orig_width = comp.get("width", 2.0)
            orig_height = comp.get("height", 0.8)

            # 스케일 및 오프셋 적용
            x = orig_x * scale + offset_x
            y = orig_y * scale + offset_y
            width = orig_width * scale
            height = orig_height * scale

            # 최소 크기 보장
            width = max(width, 0.8)
            height = max(height, 0.4)

            # 중요도(priority) 기반 색상 지원: high, medium, low
            priority = comp.get("priority", "").lower()
            priority_colors = {
                "high": self.design.BRAND_COLORS["primary"],      # 진한 네이비 (가장 중요)
                "critical": self.design.BRAND_COLORS["primary"],
                "medium": self.design.BRAND_COLORS["secondary"],  # 파랑 (중간)
                "normal": self.design.BRAND_COLORS["secondary"],
                "low": self.design.BRAND_COLORS["accent"],        # 중간 파랑 (낮음)
                "optional": RGBColor(150, 180, 200),              # 연한 파랑 (선택적)
            }

            color_key = comp.get("color", "")
            if priority and priority in priority_colors and not color_key:
                # priority 설정이 있고 color가 없으면 priority 색상 사용
                fill_color = priority_colors[priority]
            elif isinstance(color_key, str) and color_key in colors:
                fill_color = colors[color_key]
            elif isinstance(color_key, (list, tuple)) and len(color_key) == 3:
                fill_color = RGBColor(*color_key)
            else:
                fill_color = colors["primary"]

            text_color_key = comp.get("text_color", "text")
            if text_color_key in colors:
                text_color = colors[text_color_key]
            else:
                text_color = colors["text"]

            shape_type_str = comp.get("shape", "rounded_rectangle")
            shape_type = shape_types.get(shape_type_str, MSO_SHAPE.ROUNDED_RECTANGLE)

            # 폰트 크기도 스케일에 맞게 조정 (최소 10pt 보장)
            base_font_size = comp.get("font_size", 11)
            adjusted_font_size = max(10, int(base_font_size * scale))

            # 박스에는 이름만 표시 (간결하게)
            comp_name = comp.get("name", comp.get("text", ""))
            comp_desc = comp.get("description", "")

            self._add_box(
                slide,
                text=comp_name,
                x=x, y=y,
                width=width, height=height,
                fill_color=fill_color,
                text_color=text_color,
                font_size=adjusted_font_size,
                bold=comp.get("bold", True),
                shape_type=shape_type,
                border_width=comp.get("border_width", 0)
            )

            # 설명이 있으면 박스 아래에 별도 텍스트로 표시
            if comp_desc:
                desc_box = slide.shapes.add_textbox(
                    Inches(x), Inches(y + height + 0.02),
                    Inches(width), Inches(0.25)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                tf.anchor = MSO_ANCHOR.TOP
                p = tf.paragraphs[0]
                p.text = comp_desc
                p.alignment = PP_ALIGN.CENTER
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(8)
                p.font.color.rgb = RGBColor(100, 100, 100)

            component_positions[comp_id] = {
                "x": x, "y": y,
                "width": width, "height": height,
                "center_x": x + width / 2,
                "center_y": y + height / 2,
            }

        # 연결선 (자연스러운 라우팅)
        if connections:
            line_color = colors.get("accent", RGBColor(89, 89, 89))

            # 각 연결을 개별적으로 처리 (단순하고 명확한 연결선)
            for conn in connections:
                from_id = conn.get("from")
                to_id = conn.get("to")

                if from_id not in component_positions or to_id not in component_positions:
                    continue

                from_pos = component_positions[from_id]
                to_pos = component_positions[to_id]

                # 두 박스 간의 상대 위치 계산
                dx = to_pos["center_x"] - from_pos["center_x"]
                dy = to_pos["center_y"] - from_pos["center_y"]

                # 연결 방향 결정 (주요 이동 방향 기준)
                if abs(dx) > abs(dy):
                    # 주로 수평 이동
                    if dx > 0:
                        # 오른쪽으로 이동: 소스 오른쪽 → 타겟 왼쪽
                        start_x = from_pos["x"] + from_pos["width"]
                        start_y = from_pos["center_y"]
                        end_x = to_pos["x"]
                        end_y = to_pos["center_y"]
                    else:
                        # 왼쪽으로 이동: 소스 왼쪽 → 타겟 오른쪽
                        start_x = from_pos["x"]
                        start_y = from_pos["center_y"]
                        end_x = to_pos["x"] + to_pos["width"]
                        end_y = to_pos["center_y"]

                    # 수평 정렬 확인
                    if abs(start_y - end_y) < 0.15:
                        # 거의 수평: 직선 연결
                        self._add_arrow(slide, start_x, start_y, end_x, end_y,
                                       color=line_color, line_width=1.5)
                    else:
                        # 엘보우 연결 (수평 → 수직 → 수평)
                        mid_x = (start_x + end_x) / 2
                        self._add_line(slide, start_x, start_y, mid_x, start_y,
                                      color=line_color, line_width=1.5)
                        self._add_line(slide, mid_x, start_y, mid_x, end_y,
                                      color=line_color, line_width=1.5)
                        self._add_arrow(slide, mid_x, end_y, end_x, end_y,
                                       color=line_color, line_width=1.5)
                else:
                    # 주로 수직 이동
                    if dy > 0:
                        # 아래로 이동: 소스 아래 → 타겟 위
                        start_x = from_pos["center_x"]
                        start_y = from_pos["y"] + from_pos["height"]
                        end_x = to_pos["center_x"]
                        end_y = to_pos["y"]
                    else:
                        # 위로 이동: 소스 위 → 타겟 아래
                        start_x = from_pos["center_x"]
                        start_y = from_pos["y"]
                        end_x = to_pos["center_x"]
                        end_y = to_pos["y"] + to_pos["height"]

                    # 수직 정렬 확인
                    if abs(start_x - end_x) < 0.15:
                        # 거의 수직: 직선 연결
                        self._add_arrow(slide, start_x, start_y, end_x, end_y,
                                       color=line_color, line_width=1.5)
                    else:
                        # 엘보우 연결 (수직 → 수평 → 수직)
                        mid_y = (start_y + end_y) / 2
                        self._add_line(slide, start_x, start_y, start_x, mid_y,
                                      color=line_color, line_width=1.5)
                        self._add_line(slide, start_x, mid_y, end_x, mid_y,
                                      color=line_color, line_width=1.5)
                        self._add_arrow(slide, end_x, mid_y, end_x, end_y,
                                       color=line_color, line_width=1.5)

        # 레이블 (스케일 및 오프셋 적용)
        if labels:
            for label in labels:
                label_color = colors.get(label.get("color", "dark_text"), colors["dark_text"])
                orig_label_x = label.get("x", 0.5)
                orig_label_y = label.get("y", 1.0)

                # 스케일 및 오프셋 적용
                label_x = orig_label_x * scale + offset_x
                label_y = orig_label_y * scale + offset_y

                # 폰트 크기 스케일 조정
                base_label_size = label.get("font_size", 10)
                adjusted_label_size = max(8, int(base_label_size * scale))

                self._add_label(
                    slide,
                    text=label.get("text", ""),
                    x=label_x,
                    y=label_y,
                    font_size=adjusted_label_size,
                    color=label_color,
                    bold=label.get("bold", False)
                )

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def add_flowchart_slide(
        self,
        title: str,
        flow_type: str = "horizontal",
        steps: List[Dict[str, Any]] = None,
        palette: str = "blue"
    ):
        """플로우차트 슬라이드"""
        components = []
        connections = []

        if flow_type == "horizontal":
            start_x = 0.3
            y = 3.2
            box_width = 1.7
            box_height = 0.9
            gap = 0.4

            for i, step in enumerate(steps or []):
                x = start_x + i * (box_width + gap)
                comp_id = f"step_{i}"

                color_keys = ["primary", "secondary", "accent"]
                color = step.get("color", color_keys[i % len(color_keys)])

                # title과 description을 name으로 조합
                step_title = step.get("title", step.get("text", f"Step {i+1}"))
                step_desc = step.get("description", "")

                components.append({
                    "id": comp_id,
                    "name": step_title,
                    "description": step_desc,
                    "x": x, "y": y,
                    "width": box_width,
                    "height": box_height,
                    "color": color,
                    "shape": step.get("shape", "rounded_rectangle"),
                    "font_size": step.get("font_size", 12)
                })

                if i > 0:
                    connections.append({
                        "from": f"step_{i-1}",
                        "to": comp_id,
                        "direction": "right"
                    })
        else:
            x = 3.5
            start_y = 1.2
            box_width = 3.5
            box_height = 0.7
            gap = 0.4

            for i, step in enumerate(steps or []):
                y = start_y + i * (box_height + gap)
                comp_id = f"step_{i}"

                color_keys = ["primary", "secondary", "accent"]
                color = step.get("color", color_keys[i % len(color_keys)])

                # title과 description을 name으로 조합
                step_title = step.get("title", step.get("text", f"Step {i+1}"))
                step_desc = step.get("description", "")

                components.append({
                    "id": comp_id,
                    "name": step_title,
                    "description": step_desc,
                    "x": x, "y": y,
                    "width": box_width,
                    "height": box_height,
                    "color": color,
                    "shape": step.get("shape", "rounded_rectangle"),
                    "font_size": step.get("font_size", 12)
                })

                if i > 0:
                    connections.append({
                        "from": f"step_{i-1}",
                        "to": comp_id,
                        "direction": "down"
                    })

        return self.add_architecture_slide(
            title=title,
            components=components,
            connections=connections,
            palette=palette
        )

    def add_image_slide(
        self,
        title: str,
        image_path: str,
        caption: str = "",
        image_width: float = None,
        image_height: float = None
    ):
        """이미지 슬라이드 추가

        Args:
            title: 슬라이드 제목
            image_path: 이미지 파일 경로
            caption: 이미지 하단 캡션 (선택)
            image_width: 이미지 너비 (인치, 선택)
            image_height: 이미지 높이 (인치, 선택)
        """
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 제목
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 이미지 삽입
        image_file = Path(image_path)
        if image_file.exists():
            # 이미지 크기 계산
            if image_width and image_height:
                img_width = Inches(image_width)
                img_height = Inches(image_height)
            elif image_width:
                img_width = Inches(image_width)
                img_height = None
            elif image_height:
                img_width = None
                img_height = Inches(image_height)
            else:
                # 기본 크기 (슬라이드에 맞춤)
                img_width = Inches(8.0)
                img_height = None

            # 중앙 정렬
            if img_width:
                width_inches = img_width.inches if hasattr(img_width, 'inches') else img_width / 914400
                left = Inches((10.8 - width_inches) / 2)
            else:
                left = Inches(1.4)

            top = Inches(1.3)

            # 이미지 추가
            if img_height:
                slide.shapes.add_picture(
                    str(image_file), left, top,
                    width=img_width, height=img_height
                )
            else:
                slide.shapes.add_picture(
                    str(image_file), left, top,
                    width=img_width
                )

        # 캡션 추가
        if caption:
            caption_box = slide.shapes.add_textbox(
                Inches(0.4), Inches(6.3),
                Inches(10.0), Inches(0.5)
            )
            tf = caption_box.text_frame
            p = tf.paragraphs[0]
            p.text = caption
            p.alignment = PP_ALIGN.CENTER
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(11)
            p.font.italic = True
            p.font.color.rgb = RGBColor(96, 96, 96)

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def add_timeline_slide(
        self,
        title: str,
        milestones: List[Dict[str, Any]],
        style: str = "horizontal"
    ):
        """타임라인/로드맵 슬라이드

        Args:
            title: 슬라이드 제목
            milestones: 마일스톤 목록 [{date, title, description, status}]
                - date: 날짜/기간 문자열
                - title: 마일스톤 제목
                - description: 설명 (선택)
                - status: completed, current, upcoming (선택)
            style: horizontal 또는 vertical
        """
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 제목
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        if style == "horizontal":
            self._draw_horizontal_timeline(slide, milestones)
        else:
            self._draw_vertical_timeline(slide, milestones)

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def _draw_horizontal_timeline(self, slide, milestones: List[Dict[str, Any]]):
        """가로형 타임라인 그리기 - 도형 기반 개선"""
        num_items = len(milestones)
        if num_items == 0:
            return

        # 레이아웃 계산
        start_x = 0.4
        content_width = 10.0
        box_width = min(2.0, (content_width - 0.3 * (num_items - 1)) / num_items)
        gap = 0.3
        box_height = 2.8
        box_y = 1.8

        # 색상 순환 (파란색 계열 우선)
        color_keys = ["primary", "secondary", "accent", "success", "warning"]

        for i, milestone in enumerate(milestones):
            x = start_x + i * (box_width + gap)
            color_key = color_keys[i % len(color_keys)]
            box_color = self.design.BRAND_COLORS[color_key]

            # 배경 박스 (그림자)
            shadow = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.03), Inches(box_y + 0.03),
                Inches(box_width), Inches(box_height)
            )
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(200, 200, 200)
            shadow.line.fill.background()

            # 메인 박스
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(box_y),
                Inches(box_width), Inches(box_height)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = self.design.BRAND_COLORS["white"]
            box.line.color.rgb = box_color
            box.line.width = Pt(2)

            # 상단 컬러 바
            color_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(box_y),
                Inches(box_width), Inches(0.15)
            )
            color_bar.fill.solid()
            color_bar.fill.fore_color.rgb = box_color
            color_bar.line.fill.background()

            # 스텝 번호 원형
            circle_size = 0.5
            circle_x = x + (box_width - circle_size) / 2
            circle_y = box_y + 0.3
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(circle_x), Inches(circle_y),
                Inches(circle_size), Inches(circle_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = box_color
            circle.line.fill.background()

            # 스텝 번호 텍스트
            tf = circle.text_frame
            tf.anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.alignment = PP_ALIGN.CENTER
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.design.BRAND_COLORS["white"]

            # 날짜/라벨 (원 아래)
            date_text = milestone.get("date", "")
            if date_text:
                date_box = slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(circle_y + circle_size + 0.1),
                    Inches(box_width - 0.2), Inches(0.4)
                )
                tf = date_box.text_frame
                p = tf.paragraphs[0]
                p.text = date_text
                p.alignment = PP_ALIGN.CENTER
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = box_color

            # 제목
            title_text = milestone.get("title", "")
            title_y = circle_y + circle_size + 0.5
            if title_text:
                title_box = slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(title_y),
                    Inches(box_width - 0.2), Inches(0.6)
                )
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title_text
                p.alignment = PP_ALIGN.CENTER
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = self.design.BRAND_COLORS["black"]

            # 설명
            desc_text = milestone.get("description", "")
            if desc_text:
                desc_box = slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(title_y + 0.55),
                    Inches(box_width - 0.2), Inches(1.0)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = desc_text
                p.alignment = PP_ALIGN.CENTER
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(80, 80, 80)

            # 화살표 연결 (마지막 아이템 제외)
            if i < num_items - 1:
                arrow_start_x = x + box_width + 0.05
                arrow_end_x = x + box_width + gap - 0.05
                arrow_y = box_y + box_height / 2
                self._add_arrow(
                    slide,
                    arrow_start_x, arrow_y,
                    arrow_end_x, arrow_y,
                    color=self.design.BRAND_COLORS["primary"],
                    line_width=2.5
                )

    def _draw_vertical_timeline(self, slide, milestones: List[Dict[str, Any]]):
        """세로형 타임라인 그리기"""
        num_items = len(milestones)
        if num_items == 0:
            return

        # 타임라인 라인
        line_x = 2.0
        line_start_y = 1.2
        line_end_y = 6.5
        line_length = line_end_y - line_start_y

        # 메인 라인
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(line_x), Inches(line_start_y),
            Inches(0.06), Inches(line_length)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.design.BRAND_COLORS["primary"]
        line.line.fill.background()

        # 각 마일스톤 배치
        spacing = line_length / max(num_items, 1)

        status_colors = {
            "completed": self.design.BRAND_COLORS["success"],
            "current": self.design.BRAND_COLORS["warning"],
            "upcoming": RGBColor(180, 180, 180),
        }

        for i, milestone in enumerate(milestones):
            y = line_start_y + i * spacing + spacing / 2

            status = milestone.get("status", "upcoming")
            dot_color = status_colors.get(status, status_colors["upcoming"])

            # 동그라미 (마일스톤 포인트)
            dot_size = 0.3
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(line_x - dot_size/2 + 0.03), Inches(y - dot_size/2),
                Inches(dot_size), Inches(dot_size)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = dot_color
            dot.line.color.rgb = self.design.BRAND_COLORS["white"]
            dot.line.width = Pt(2)

            # 날짜 (왼쪽)
            date_text = milestone.get("date", "")
            if date_text:
                date_box = slide.shapes.add_textbox(
                    Inches(0.3), Inches(y - 0.15),
                    Inches(1.5), Inches(0.4)
                )
                tf = date_box.text_frame
                p = tf.paragraphs[0]
                p.text = date_text
                p.alignment = PP_ALIGN.RIGHT
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = self.design.BRAND_COLORS["primary"]

            # 제목 및 설명 카드 (오른쪽)
            card_x = 2.6
            card_width = 7.0
            card_height = spacing * 0.8

            # 카드 배경
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(card_x), Inches(y - card_height/2),
                Inches(card_width), Inches(card_height)
            )

            # 상태에 따른 배경색
            if status == "completed":
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(240, 255, 240)
                card.line.color.rgb = self.design.BRAND_COLORS["success"]
            elif status == "current":
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(255, 250, 230)
                card.line.color.rgb = self.design.BRAND_COLORS["warning"]
            else:
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(248, 248, 248)
                card.line.color.rgb = RGBColor(200, 200, 200)
            card.line.width = Pt(1)

            # 제목
            title_text = milestone.get("title", "")
            if title_text:
                title_box = slide.shapes.add_textbox(
                    Inches(card_x + 0.15), Inches(y - card_height/2 + 0.1),
                    Inches(card_width - 0.3), Inches(0.35)
                )
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = title_text
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = self.design.BRAND_COLORS["black"]

            # 설명
            desc_text = milestone.get("description", "")
            if desc_text:
                desc_box = slide.shapes.add_textbox(
                    Inches(card_x + 0.15), Inches(y - card_height/2 + 0.4),
                    Inches(card_width - 0.3), Inches(card_height - 0.5)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = desc_text
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(80, 80, 80)

    def add_stats_slide(
        self,
        title: str,
        stats: List[Dict[str, Any]] = None,
        style: str = "cards"
    ):
        """통계/수치 강조 슬라이드

        Args:
            title: 슬라이드 제목
            stats: 통계 목록 [{value, label, unit, description, color, icon}]
                - value: 수치 값 (필수)
                - label: 레이블 (필수)
                - unit: 단위 (%, 건, 명 등)
                - description: 부가 설명
                - color: 색상 (primary, success, warning 등)
                - icon: 아이콘 이모지
            style: cards 또는 inline
        """
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 제목
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        stats = stats or []
        num_stats = len(stats)
        if num_stats == 0:
            self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
            return slide

        if style == "cards":
            self._draw_stats_cards(slide, stats)
        else:
            self._draw_stats_inline(slide, stats)

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def _draw_stats_cards(self, slide, stats: List[Dict[str, Any]]):
        """카드 형태의 통계 표시"""
        num_stats = len(stats)

        # 레이아웃 계산 (최대 4개까지 한 줄)
        if num_stats <= 4:
            cols = num_stats
            rows = 1
        else:
            cols = min(4, num_stats)
            rows = (num_stats + cols - 1) // cols

        card_width = (10.0 - (cols - 1) * 0.3) / cols
        card_height = 2.2 if rows == 1 else 1.8
        start_x = 0.4
        start_y = 1.5

        color_keys = ["primary", "secondary", "accent", "success", "warning"]

        for i, stat in enumerate(stats):
            col = i % cols
            row = i // cols
            x = start_x + col * (card_width + 0.3)
            y = start_y + row * (card_height + 0.3)

            # 색상 결정
            color_key = stat.get("color", color_keys[i % len(color_keys)])
            if color_key in self.design.BRAND_COLORS:
                accent_color = self.design.BRAND_COLORS[color_key]
            else:
                accent_color = self.design.BRAND_COLORS["primary"]

            # 카드 배경
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y),
                Inches(card_width), Inches(card_height)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.design.BRAND_COLORS["white"]
            card.line.color.rgb = RGBColor(230, 230, 230)
            card.line.width = Pt(1)

            # 상단 악센트 바
            accent_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y),
                Inches(card_width), Inches(0.1)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = accent_color
            accent_bar.line.fill.background()

            # 아이콘 (있는 경우)
            icon = stat.get("icon", "")
            if icon:
                icon_box = slide.shapes.add_textbox(
                    Inches(x + 0.15), Inches(y + 0.25),
                    Inches(0.5), Inches(0.5)
                )
                tf = icon_box.text_frame
                p = tf.paragraphs[0]
                p.text = icon
                p.font.size = Pt(24)

            # 수치
            value = str(stat.get("value", "0"))
            unit = stat.get("unit", "")
            value_text = f"{value}{unit}"

            value_x = x + 0.15 if not icon else x + 0.6
            value_box = slide.shapes.add_textbox(
                Inches(value_x), Inches(y + 0.3),
                Inches(card_width - 0.3), Inches(0.8)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value_text
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = accent_color

            # 레이블
            label = stat.get("label", "")
            label_box = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(y + 1.1),
                Inches(card_width - 0.3), Inches(0.4)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = self.design.BRAND_COLORS["black"]

            # 설명 (있는 경우)
            description = stat.get("description", "")
            if description and card_height > 1.8:
                desc_box = slide.shapes.add_textbox(
                    Inches(x + 0.15), Inches(y + 1.5),
                    Inches(card_width - 0.3), Inches(0.5)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = description
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(120, 120, 120)

    def _draw_stats_inline(self, slide, stats: List[Dict[str, Any]]):
        """인라인 형태의 통계 표시"""
        num_stats = len(stats)
        stat_width = 10.0 / num_stats
        start_y = 2.5

        color_keys = ["primary", "secondary", "accent", "success", "warning"]

        for i, stat in enumerate(stats):
            x = 0.4 + i * stat_width
            center_x = x + stat_width / 2

            # 색상
            color_key = stat.get("color", color_keys[i % len(color_keys)])
            if color_key in self.design.BRAND_COLORS:
                accent_color = self.design.BRAND_COLORS[color_key]
            else:
                accent_color = self.design.BRAND_COLORS["primary"]

            # 원형 배경
            circle_size = 1.8
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(center_x - circle_size/2), Inches(start_y),
                Inches(circle_size), Inches(circle_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = accent_color
            circle.line.fill.background()

            # 수치
            value = str(stat.get("value", "0"))
            unit = stat.get("unit", "")
            value_text = f"{value}{unit}"

            value_box = slide.shapes.add_textbox(
                Inches(center_x - 0.9), Inches(start_y + 0.5),
                Inches(1.8), Inches(0.8)
            )
            tf = value_box.text_frame
            p = tf.paragraphs[0]
            p.text = value_text
            p.alignment = PP_ALIGN.CENTER
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = self.design.BRAND_COLORS["white"]

            # 레이블
            label = stat.get("label", "")
            label_box = slide.shapes.add_textbox(
                Inches(x), Inches(start_y + circle_size + 0.3),
                Inches(stat_width), Inches(0.5)
            )
            tf = label_box.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.alignment = PP_ALIGN.CENTER
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.design.BRAND_COLORS["black"]

            # 설명
            description = stat.get("description", "")
            if description:
                desc_box = slide.shapes.add_textbox(
                    Inches(x), Inches(start_y + circle_size + 0.7),
                    Inches(stat_width), Inches(0.5)
                )
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = description
                p.alignment = PP_ALIGN.CENTER
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(100, 100, 100)

    def add_two_column_slide(
        self,
        title: str,
        left_content: Dict[str, Any],
        right_content: Dict[str, Any],
        column_ratio: str = "50:50"
    ):
        """2단 레이아웃 슬라이드

        Args:
            title: 슬라이드 제목
            left_content: 왼쪽 컬럼 내용 {type, title, items/text/image_path}
            right_content: 오른쪽 컬럼 내용 {type, title, items/text/image_path}
            column_ratio: 컬럼 비율 (50:50, 40:60, 60:40, 30:70, 70:30)
        """
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 제목
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 컬럼 비율 계산
        ratio_map = {
            "50:50": (0.5, 0.5),
            "40:60": (0.4, 0.6),
            "60:40": (0.6, 0.4),
            "30:70": (0.3, 0.7),
            "70:30": (0.7, 0.3),
        }
        left_ratio, right_ratio = ratio_map.get(column_ratio, (0.5, 0.5))

        total_width = 9.8
        gap = 0.3
        left_width = total_width * left_ratio - gap / 2
        right_width = total_width * right_ratio - gap / 2

        left_x = 0.5
        right_x = left_x + left_width + gap
        content_y = 1.2
        content_height = 5.5

        self._draw_column_content(
            slide, left_content,
            left_x, content_y, left_width, content_height
        )

        self._draw_column_content(
            slide, right_content,
            right_x, content_y, right_width, content_height
        )

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def add_tree_slide(
        self,
        title: str,
        tree_structure: List[Dict[str, Any]],
        descriptions: Dict[str, str] = None
    ):
        """디렉토리 트리 구조 슬라이드 (Ubuntu tree 명령어 스타일)

        Args:
            title: 슬라이드 제목
            tree_structure: 트리 구조 리스트
                [{name: "폴더명", children: [...], description: "설명"}]
            descriptions: 폴더별 설명 딕셔너리 (선택)
        """
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 제목
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 트리 영역과 설명 영역 분리 (40:60 비율)
        tree_x = 0.5
        tree_width = 3.8
        desc_x = 4.5
        desc_width = 5.8
        content_y = 1.2
        content_height = 5.5

        # 트리 배경 박스
        tree_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(tree_x), Inches(content_y),
            Inches(tree_width), Inches(content_height)
        )
        tree_bg.fill.solid()
        tree_bg.fill.fore_color.rgb = RGBColor(245, 248, 250)
        tree_bg.line.color.rgb = self.design.BRAND_COLORS["secondary"]
        tree_bg.line.width = Pt(1)

        # 트리 텍스트 생성
        tree_lines = []
        self._build_tree_lines(tree_structure, tree_lines, "", True)

        # 트리 텍스트 박스
        tree_text = slide.shapes.add_textbox(
            Inches(tree_x + 0.2), Inches(content_y + 0.15),
            Inches(tree_width - 0.4), Inches(content_height - 0.3)
        )
        tf = tree_text.text_frame
        tf.word_wrap = False

        for i, line in enumerate(tree_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.name = "Consolas"  # 고정폭 폰트
            p.font.size = Pt(11)
            p.font.color.rgb = self.design.BRAND_COLORS["dark"]
            p.line_spacing = 1.2

        # 설명 영역 - 폴더별 설명 표시
        desc_data = descriptions or {}
        # tree_structure에서 description 추출
        self._extract_descriptions(tree_structure, desc_data)

        if desc_data:
            # 설명 제목
            desc_title = slide.shapes.add_textbox(
                Inches(desc_x), Inches(content_y),
                Inches(desc_width), Inches(0.4)
            )
            tf = desc_title.text_frame
            p = tf.paragraphs[0]
            p.text = "폴더 설명"
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = self.design.BRAND_COLORS["primary"]

            # 각 폴더 설명을 카드 형태로 표시
            card_y = content_y + 0.5
            card_height = 0.6
            card_gap = 0.1

            for folder_name, desc in list(desc_data.items())[:8]:  # 최대 8개
                if card_y + card_height > content_y + content_height:
                    break

                # 폴더명 배경
                folder_bg = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(desc_x), Inches(card_y),
                    Inches(desc_width), Inches(card_height)
                )
                folder_bg.fill.solid()
                folder_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
                folder_bg.line.color.rgb = RGBColor(220, 220, 220)
                folder_bg.line.width = Pt(1)

                # 좌측 컬러바
                color_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(desc_x), Inches(card_y),
                    Inches(0.06), Inches(card_height)
                )
                color_bar.fill.solid()
                color_bar.fill.fore_color.rgb = self.design.BRAND_COLORS["secondary"]
                color_bar.line.fill.background()

                # 폴더명 (굵게)
                folder_text = slide.shapes.add_textbox(
                    Inches(desc_x + 0.15), Inches(card_y + 0.08),
                    Inches(desc_width - 0.3), Inches(0.25)
                )
                tf = folder_text.text_frame
                tf.word_wrap = False
                tf.anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = f"📁 {folder_name}/"
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = self.design.BRAND_COLORS["primary"]

                # 설명
                desc_text = slide.shapes.add_textbox(
                    Inches(desc_x + 0.15), Inches(card_y + 0.32),
                    Inches(desc_width - 0.3), Inches(0.25)
                )
                tf = desc_text.text_frame
                tf.word_wrap = True
                tf.anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = desc
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(80, 80, 80)

                card_y += card_height + card_gap

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def _build_tree_lines(
        self,
        items: List[Dict[str, Any]],
        lines: List[str],
        prefix: str = "",
        is_root: bool = False
    ):
        """트리 라인 생성 (재귀)"""
        for i, item in enumerate(items):
            is_last = (i == len(items) - 1)
            name = item.get("name", "")

            if is_root and i == 0 and not prefix:
                # 루트 레벨
                lines.append(f"📦 {name}/")
                child_prefix = "   "
            else:
                connector = "└── " if is_last else "├── "
                icon = "📁" if item.get("children") else "📄"
                suffix = "/" if item.get("children") else ""
                lines.append(f"{prefix}{connector}{icon} {name}{suffix}")
                child_prefix = prefix + ("    " if is_last else "│   ")

            # 자식 폴더 처리
            children = item.get("children", [])
            if children:
                self._build_tree_lines(children, lines, child_prefix, False)

    def _extract_descriptions(
        self,
        items: List[Dict[str, Any]],
        desc_dict: Dict[str, str]
    ):
        """트리 구조에서 설명 추출 (재귀)"""
        for item in items:
            name = item.get("name", "")
            desc = item.get("description", "")
            if name and desc and name not in desc_dict:
                desc_dict[name] = desc
            children = item.get("children", [])
            if children:
                self._extract_descriptions(children, desc_dict)

    def _draw_column_content(
        self,
        slide,
        content: Dict[str, Any],
        x: float,
        y: float,
        width: float,
        height: float
    ):
        """컬럼 내용 그리기"""
        content_type = content.get("type", "bullets")
        col_title = content.get("title", "")

        current_y = y

        # 컬럼 제목 (있는 경우)
        if col_title:
            # 제목 배경
            title_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(current_y),
                Inches(width), Inches(0.45)
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = self.design.BRAND_COLORS["primary"]
            title_bg.line.fill.background()

            title_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(current_y + 0.08),
                Inches(width - 0.2), Inches(0.35)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = col_title
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.design.BRAND_COLORS["white"]

            current_y += 0.55

        # 컬럼 배경 (선택적)
        if content.get("background", False):
            bg_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(current_y),
                Inches(width), Inches(height - (current_y - y))
            )
            bg_color = content.get("bg_color", "light")
            if bg_color in self.design.BRAND_COLORS:
                bg_box.fill.solid()
                bg_box.fill.fore_color.rgb = self.design.BRAND_COLORS[bg_color]
            else:
                bg_box.fill.solid()
                bg_box.fill.fore_color.rgb = RGBColor(250, 250, 250)
            bg_box.line.fill.background()

        remaining_height = height - (current_y - y)

        if content_type == "bullets":
            # 불릿 포인트 목록
            items = content.get("items", [])
            for i, item in enumerate(items):
                item_y = current_y + i * 0.55
                item_box = slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(item_y),
                    Inches(width - 0.2), Inches(0.5)
                )
                tf = item_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"• {item}"
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(12)
                p.font.color.rgb = self.design.BRAND_COLORS["black"]

        elif content_type == "text":
            # 자유 텍스트
            text = content.get("text", "")
            text_box = slide.shapes.add_textbox(
                Inches(x + 0.1), Inches(current_y + 0.1),
                Inches(width - 0.2), Inches(remaining_height - 0.2)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(12)
            p.font.color.rgb = self.design.BRAND_COLORS["black"]

        elif content_type == "image":
            # 이미지
            image_path = content.get("image_path", "")
            if image_path and Path(image_path).exists():
                try:
                    slide.shapes.add_picture(
                        image_path,
                        Inches(x + 0.1), Inches(current_y + 0.1),
                        width=Inches(width - 0.2)
                    )
                except Exception:
                    # 이미지 로드 실패 시 플레이스홀더
                    placeholder = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(x + 0.1), Inches(current_y + 0.1),
                        Inches(width - 0.2), Inches(remaining_height - 0.2)
                    )
                    placeholder.fill.solid()
                    placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)

        elif content_type == "numbered":
            # 번호 목록
            items = content.get("items", [])
            for i, item in enumerate(items):
                item_y = current_y + i * 0.55
                item_box = slide.shapes.add_textbox(
                    Inches(x + 0.1), Inches(item_y),
                    Inches(width - 0.2), Inches(0.5)
                )
                tf = item_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"{i + 1}. {item}"
                p.font.name = self.design.FONT_BODY
                p.font.size = Pt(12)
                p.font.color.rgb = self.design.BRAND_COLORS["black"]

    def set_footer(
        self,
        text: str = None,
        show_slide_number: bool = True,
        show_date: bool = False
    ):
        """푸터 설정 (이후 생성되는 슬라이드에 적용)

        Args:
            text: 푸터 텍스트 (None이면 기본 템플릿 사용)
            show_slide_number: 슬라이드 번호 표시 여부
            show_date: 날짜 표시 여부
        """
        self._footer_text = text
        self._show_slide_number = show_slide_number
        self._show_footer_date = show_date

    def add_footer_to_slide(self, slide, text: str = None):
        """특정 슬라이드에 푸터 추가

        Args:
            slide: 슬라이드 객체
            text: 푸터 텍스트
        """
        if text:
            footer_box = slide.shapes.add_textbox(
                Inches(0.4), Inches(7.0),
                Inches(9.0), Inches(0.3)
            )
            tf = footer_box.text_frame
            p = tf.paragraphs[0]
            p.text = text
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(128, 128, 128)

    def add_chart_slide(
        self,
        title: str,
        chart_type: str,
        categories: List[str],
        series: List[Dict[str, Any]],
        chart_title: str = None,
        show_legend: bool = True,
        legend_position: str = "bottom"
    ):
        """차트 슬라이드 추가

        Args:
            title: 슬라이드 제목
            chart_type: 차트 유형 (bar, column, line, pie, doughnut, area, radar)
            categories: 카테고리 목록 (X축 레이블)
            series: 데이터 시리즈 목록 [{name, values}]
            chart_title: 차트 제목 (선택)
            show_legend: 범례 표시 여부
            legend_position: 범례 위치 (top, bottom, left, right)
        """
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 슬라이드 제목
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 차트 타입 매핑
        chart_type_map = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
            "line": XL_CHART_TYPE.LINE,
            "line_markers": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "area": XL_CHART_TYPE.AREA,
            "area_stacked": XL_CHART_TYPE.AREA_STACKED,
            "radar": XL_CHART_TYPE.RADAR,
            "radar_filled": XL_CHART_TYPE.RADAR_FILLED,
        }

        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        # 차트 데이터 생성
        chart_data = CategoryChartData()
        chart_data.categories = categories

        for s in series:
            chart_data.add_series(s.get("name", "Series"), s.get("values", []))

        # 차트 위치 및 크기
        x, y, cx, cy = Inches(0.5), Inches(1.3), Inches(9.8), Inches(5.5)

        # 차트 추가
        chart = slide.shapes.add_chart(
            xl_chart_type, x, y, cx, cy, chart_data
        ).chart

        # 차트 제목 설정
        if chart_title:
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = chart_title
            chart.chart_title.text_frame.paragraphs[0].font.name = "맑은 고딕"
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
            chart.chart_title.text_frame.paragraphs[0].font.bold = True

        # 범례 설정
        if show_legend and chart_type not in ["pie", "doughnut"]:
            chart.has_legend = True
            legend_pos_map = {
                "top": XL_LEGEND_POSITION.TOP,
                "bottom": XL_LEGEND_POSITION.BOTTOM,
                "left": XL_LEGEND_POSITION.LEFT,
                "right": XL_LEGEND_POSITION.RIGHT,
            }
            chart.legend.position = legend_pos_map.get(
                legend_position, XL_LEGEND_POSITION.BOTTOM
            )
            chart.legend.include_in_layout = False

        # 시리즈 색상 설정
        color_keys = ["primary", "secondary", "accent", "success", "warning", "danger"]
        if hasattr(chart, 'series'):
            for i, s in enumerate(chart.series):
                color_key = color_keys[i % len(color_keys)]
                fill = s.format.fill
                fill.solid()
                fill.fore_color.rgb = self.design.BRAND_COLORS[color_key]

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def add_org_chart_slide(
        self,
        title: str,
        org_data: Dict[str, Any],
        style: str = "horizontal"
    ):
        """조직도 슬라이드 추가

        Args:
            title: 슬라이드 제목
            org_data: 조직 구조 데이터 (계층적 딕셔너리)
                {
                    "name": "CEO",
                    "title": "대표이사",
                    "children": [
                        {"name": "CTO", "title": "기술이사", "children": [...]},
                        {"name": "CFO", "title": "재무이사", "children": [...]}
                    ]
                }
            style: horizontal (가로) 또는 vertical (세로)
        """
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 슬라이드 제목
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 조직도 그리기
        if style == "horizontal":
            self._draw_org_chart_horizontal(slide, org_data)
        else:
            self._draw_org_chart_vertical(slide, org_data)

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def _draw_org_chart_horizontal(self, slide, org_data: Dict[str, Any]):
        """가로형 조직도 그리기 (위에서 아래로)"""
        # 전체 구조 분석 - 각 레벨별 노드 수 계산
        levels = self._analyze_org_structure(org_data)
        max_width = max(len(level) for level in levels)

        # 박스 크기 설정
        box_width = min(2.0, 9.0 / max_width - 0.2)
        box_height = 0.7
        v_gap = 0.5
        h_gap = 0.3

        start_y = 1.3

        # 각 레벨 그리기
        for level_idx, level_nodes in enumerate(levels):
            num_nodes = len(level_nodes)
            total_width = num_nodes * box_width + (num_nodes - 1) * h_gap
            start_x = (10.8 - total_width) / 2

            y = start_y + level_idx * (box_height + v_gap)

            for node_idx, node in enumerate(level_nodes):
                x = start_x + node_idx * (box_width + h_gap)

                # 색상 결정 (레벨별)
                color_keys = ["primary", "secondary", "accent", "success"]
                color = self.design.BRAND_COLORS[color_keys[level_idx % len(color_keys)]]

                # 박스 그리기
                self._draw_org_box(slide, node, x, y, box_width, box_height, color)

                # 노드 위치 저장 (연결선용)
                node["_x"] = x + box_width / 2
                node["_y"] = y
                node["_width"] = box_width
                node["_height"] = box_height

        # 연결선 그리기
        self._draw_org_connections(slide, org_data, "vertical")

    def _draw_org_chart_vertical(self, slide, org_data: Dict[str, Any]):
        """세로형 조직도 그리기 (왼쪽에서 오른쪽으로)"""
        levels = self._analyze_org_structure(org_data)
        max_height = max(len(level) for level in levels)

        # 박스 크기 설정
        box_width = 2.0
        box_height = min(0.8, 5.5 / max_height - 0.2)
        h_gap = 0.4
        v_gap = 0.2

        start_x = 0.5

        # 각 레벨 그리기
        for level_idx, level_nodes in enumerate(levels):
            num_nodes = len(level_nodes)
            total_height = num_nodes * box_height + (num_nodes - 1) * v_gap
            start_y = 1.3 + (5.5 - total_height) / 2

            x = start_x + level_idx * (box_width + h_gap)

            for node_idx, node in enumerate(level_nodes):
                y = start_y + node_idx * (box_height + v_gap)

                # 색상 결정
                color_keys = ["primary", "secondary", "accent", "success"]
                color = self.design.BRAND_COLORS[color_keys[level_idx % len(color_keys)]]

                # 박스 그리기
                self._draw_org_box(slide, node, x, y, box_width, box_height, color)

                # 노드 위치 저장
                node["_x"] = x
                node["_y"] = y + box_height / 2
                node["_width"] = box_width
                node["_height"] = box_height

        # 연결선 그리기
        self._draw_org_connections(slide, org_data, "horizontal")

    def _analyze_org_structure(self, org_data: Dict[str, Any]) -> List[List[Dict]]:
        """조직 구조를 레벨별로 분석"""
        levels = []

        def traverse(node, level):
            while len(levels) <= level:
                levels.append([])
            levels[level].append(node)
            for child in node.get("children", []):
                traverse(child, level + 1)

        traverse(org_data, 0)
        return levels

    def _draw_org_box(
        self,
        slide,
        node: Dict[str, Any],
        x: float,
        y: float,
        width: float,
        height: float,
        color: RGBColor
    ):
        """조직도 박스 그리기"""
        # 박스
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # 이름
        name = node.get("name", "")
        title_text = node.get("title", "")

        if title_text:
            # 이름과 직함 모두 표시
            text = f"{name}\n{title_text}"
            font_size = 9
        else:
            text = name
            font_size = 11

        tf = box.text_frame
        tf.word_wrap = True
        tf.anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = self.design.FONT_BODY
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = self.design.BRAND_COLORS["white"]

    def _draw_org_connections(self, slide, org_data: Dict[str, Any], direction: str):
        """조직도 연결선 그리기"""
        def draw_connections(parent):
            children = parent.get("children", [])
            if not children:
                return

            parent_x = parent.get("_x", 0)
            parent_y = parent.get("_y", 0)
            parent_h = parent.get("_height", 0.7)
            parent_w = parent.get("_width", 2.0)

            for child in children:
                child_x = child.get("_x", 0)
                child_y = child.get("_y", 0)
                child_h = child.get("_height", 0.7)
                child_w = child.get("_width", 2.0)

                if direction == "vertical":
                    # 위에서 아래로
                    start_x = parent_x
                    start_y = parent_y + parent_h
                    end_x = child_x
                    end_y = child_y
                else:
                    # 왼쪽에서 오른쪽으로
                    start_x = parent_x + parent_w
                    start_y = parent_y
                    end_x = child_x
                    end_y = child_y

                # 연결선 그리기
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(start_x), Inches(start_y),
                    Inches(end_x), Inches(end_y)
                )
                connector.line.color.rgb = RGBColor(150, 150, 150)
                connector.line.width = Pt(1.5)

                # 자식들의 연결선 그리기
                draw_connections(child)

        draw_connections(org_data)

    def add_summary_slide(
        self,
        title: str,
        points: List[str],
        highlight_text: str = None
    ):
        """요약 슬라이드"""
        layout = self._get_layout(self._get_content_layout_name())
        slide = self.prs.slides.add_slide(layout)

        # 제목 (진한 검정)
        title_shape = slide.shapes.title
        title_shape.text = title
        for para in title_shape.text_frame.paragraphs:
            para.font.name = self.design.FONT_TITLE
            para.font.size = Pt(24)
            para.font.bold = True
            para.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 포인트들
        start_y = 1.0
        for i, point in enumerate(points):
            y = start_y + i * 0.85

            # 체크 아이콘
            self._add_icon_box(
                slide, "✓",
                x=0.5, y=y + 0.05,
                size=0.4,
                bg_color=self.design.BRAND_COLORS["success"]
            )

            # 텍스트 (검정)
            point_box = slide.shapes.add_textbox(
                Inches(1.1), Inches(y),
                Inches(9.0), Inches(0.7)
            )
            tf = point_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = point
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(15)
            p.font.color.rgb = self.design.BRAND_COLORS["black"]

        # 강조 텍스트
        if highlight_text:
            highlight_y = start_y + len(points) * 0.85 + 0.3

            highlight_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.4), Inches(highlight_y),
                Inches(10.0), Inches(0.8)
            )
            highlight_box.fill.solid()
            highlight_box.fill.fore_color.rgb = RGBColor(255, 248, 220)
            highlight_box.line.color.rgb = self.design.BRAND_COLORS["highlight"]
            highlight_box.line.width = Pt(2)

            text_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(highlight_y + 0.2),
                Inches(9.6), Inches(0.5)
            )
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"💡 {highlight_text}"
            p.font.name = self.design.FONT_BODY
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = self.design.BRAND_COLORS["black"]

        self._clear_unused_placeholders(slide, used_placeholder_idx=[0])
        return slide

    def save(self, output_path: Optional[str] = None) -> Path:
        """PPT 파일 저장"""
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = OUTPUT_DIR / f"presentation_{timestamp}.pptx"
        else:
            output_path = Path(output_path)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        return output_path

    def export_pdf(self, output_path: Optional[str] = None, pptx_path: Optional[str] = None) -> Optional[Path]:
        """PDF로 내보내기

        LibreOffice가 설치되어 있어야 합니다.
        설치: sudo apt install libreoffice

        Args:
            output_path: PDF 출력 경로 (None이면 PPTX와 같은 이름으로 생성)
            pptx_path: 변환할 PPTX 파일 경로 (None이면 현재 프레젠테이션 저장 후 변환)

        Returns:
            생성된 PDF 파일 경로 또는 None (실패 시)
        """
        import subprocess
        import shutil

        # LibreOffice 확인
        libreoffice_path = shutil.which("libreoffice") or shutil.which("soffice")
        if not libreoffice_path:
            print("경고: LibreOffice가 설치되어 있지 않습니다.")
            print("설치 방법: sudo apt install libreoffice")
            return None

        # PPTX 파일 준비
        if pptx_path:
            pptx_file = Path(pptx_path)
        else:
            # 임시 저장
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            pptx_file = OUTPUT_DIR / f"temp_{timestamp}.pptx"
            self.prs.save(str(pptx_file))

        if not pptx_file.exists():
            print(f"오류: PPTX 파일을 찾을 수 없습니다: {pptx_file}")
            return None

        # PDF 출력 경로 설정
        if output_path:
            pdf_file = Path(output_path)
        else:
            pdf_file = pptx_file.with_suffix(".pdf")

        # LibreOffice로 변환
        try:
            result = subprocess.run(
                [
                    libreoffice_path,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", str(pdf_file.parent),
                    str(pptx_file)
                ],
                capture_output=True,
                text=True,
                timeout=120  # 2분 타임아웃
            )

            if result.returncode != 0:
                print(f"PDF 변환 오류: {result.stderr}")
                return None

            # 생성된 PDF 파일 이름 확인 (LibreOffice는 원본 이름 기반으로 생성)
            generated_pdf = pptx_file.with_suffix(".pdf")
            if generated_pdf.parent != pdf_file.parent:
                generated_pdf = pdf_file.parent / f"{pptx_file.stem}.pdf"

            # 원하는 이름으로 변경
            if generated_pdf != pdf_file and generated_pdf.exists():
                generated_pdf.rename(pdf_file)

            # 임시 PPTX 삭제 (pptx_path가 None이었던 경우)
            if not pptx_path and pptx_file.exists() and "temp_" in pptx_file.name:
                pptx_file.unlink()

            if pdf_file.exists():
                return pdf_file
            elif generated_pdf.exists():
                return generated_pdf
            else:
                print("PDF 파일 생성에 실패했습니다.")
                return None

        except subprocess.TimeoutExpired:
            print("PDF 변환 시간 초과")
            return None
        except Exception as e:
            print(f"PDF 변환 중 오류 발생: {e}")
            return None


def create_from_config(config: Dict[str, Any], output_path: Optional[str] = None, export_pdf: bool = False) -> Path:
    """설정 딕셔너리로부터 PPT 생성

    Args:
        config: 설정 딕셔너리
            - cover: 표지 정보
            - slides: 슬라이드 목록
            - settings: 전역 설정 (선택)
                - show_page_numbers: 페이지 번호 표시 여부 (기본: True)
                - theme_name: 사전 정의된 테마 (default, dark, green, purple, warm)
                - theme_path: 외부 테마 파일 경로
        output_path: 출력 파일 경로
    """
    # 전역 설정 처리
    settings = config.get("settings", {})
    show_page_numbers = settings.get("show_page_numbers", True)
    theme_name = settings.get("theme_name")
    theme_path = settings.get("theme_path")
    card_style = settings.get("card_style")  # classic, gradient, modern, solid

    generator = PPTGenerator(
        show_page_numbers=show_page_numbers,
        theme_name=theme_name,
        theme_path=theme_path
    )

    # 카드 스타일 적용 (전역 설정)
    if card_style and card_style in DesignSystem.CARD_STYLES:
        generator.design.card_style = card_style

    # 표지
    if "cover" in config:
        cover = config["cover"]
        generator.add_cover_slide(
            title=cover.get("title", "제목 없음"),
            date=cover.get("date"),
            author=cover.get("author", ""),
            report_type=cover.get("report_type", "정보공유")
        )

    # 슬라이드들
    for slide_config in config.get("slides", []):
        slide_type = slide_config.get("type", "content")

        if slide_type == "section":
            generator.add_section_slide(
                section_number=slide_config.get("number", "1"),
                section_title=slide_config.get("title", ""),
                subtitle=slide_config.get("subtitle", "")
            )
        elif slide_type == "content":
            generator.add_content_slide(
                title=slide_config.get("title", ""),
                content=slide_config.get("content", [])
            )
        elif slide_type == "content_boxed":
            generator.add_content_boxed_slide(
                title=slide_config.get("title", ""),
                sections=slide_config.get("sections", []),
                columns=slide_config.get("columns", 1)
            )
        elif slide_type == "content_icons":
            generator.add_content_with_icons(
                title=slide_config.get("title", ""),
                items=slide_config.get("items", [])
            )
        elif slide_type == "comparison":
            generator.add_comparison_slide(
                title=slide_config.get("title", ""),
                left_title=slide_config.get("left_title", "Before"),
                left_items=slide_config.get("left_items", []),
                right_title=slide_config.get("right_title", "After"),
                right_items=slide_config.get("right_items", []),
                left_color=slide_config.get("left_color", "danger"),
                right_color=slide_config.get("right_color", "success")
            )
        elif slide_type == "text":
            generator.add_text_slide(
                title=slide_config.get("title", ""),
                text_blocks=slide_config.get("text_blocks", [])
            )
        elif slide_type == "table":
            generator.add_table_slide(
                title=slide_config.get("title", ""),
                headers=slide_config.get("headers", []),
                rows=slide_config.get("rows", []),
                col_widths=slide_config.get("col_widths"),
                highlight_rows=slide_config.get("highlight_rows")
            )
        elif slide_type == "cards":
            generator.add_cards_slide(
                title=slide_config.get("title", ""),
                cards=slide_config.get("cards", []),
                columns=slide_config.get("columns", 3),
                card_style=slide_config.get("card_style")  # 슬라이드별 오버라이드 가능
            )
        elif slide_type == "architecture":
            generator.add_architecture_slide(
                title=slide_config.get("title", ""),
                components=slide_config.get("components", []),
                connections=slide_config.get("connections", []),
                labels=slide_config.get("labels", []),
                palette=slide_config.get("palette", "blue")
            )
        elif slide_type == "flowchart":
            generator.add_flowchart_slide(
                title=slide_config.get("title", ""),
                flow_type=slide_config.get("flow_type", "horizontal"),
                steps=slide_config.get("steps", []),
                palette=slide_config.get("palette", "blue")
            )
        elif slide_type == "summary":
            generator.add_summary_slide(
                title=slide_config.get("title", ""),
                points=slide_config.get("points", []),
                highlight_text=slide_config.get("highlight")
            )
        elif slide_type == "image":
            generator.add_image_slide(
                title=slide_config.get("title", ""),
                image_path=slide_config.get("image_path", ""),
                caption=slide_config.get("caption", ""),
                image_width=slide_config.get("image_width"),
                image_height=slide_config.get("image_height")
            )
        elif slide_type == "timeline":
            generator.add_timeline_slide(
                title=slide_config.get("title", ""),
                milestones=slide_config.get("milestones", []),
                style=slide_config.get("style", "horizontal")
            )
        elif slide_type == "stats":
            generator.add_stats_slide(
                title=slide_config.get("title", ""),
                stats=slide_config.get("stats", []),
                style=slide_config.get("style", "cards")
            )
        elif slide_type == "two_column":
            generator.add_two_column_slide(
                title=slide_config.get("title", ""),
                left_content=slide_config.get("left_content", slide_config.get("left", {})),
                right_content=slide_config.get("right_content", slide_config.get("right", {})),
                column_ratio=slide_config.get("column_ratio", "50:50")
            )
        elif slide_type == "chart":
            generator.add_chart_slide(
                title=slide_config.get("title", ""),
                chart_type=slide_config.get("chart_type", "column"),
                categories=slide_config.get("categories", []),
                series=slide_config.get("series", []),
                chart_title=slide_config.get("chart_title"),
                show_legend=slide_config.get("show_legend", True),
                legend_position=slide_config.get("legend_position", "bottom")
            )
        elif slide_type == "org_chart":
            generator.add_org_chart_slide(
                title=slide_config.get("title", ""),
                org_data=slide_config.get("org_data", {}),
                style=slide_config.get("style", "horizontal")
            )
        elif slide_type == "tree":
            generator.add_tree_slide(
                title=slide_config.get("title", ""),
                tree_structure=slide_config.get("tree_structure", []),
                descriptions=slide_config.get("descriptions", {})
            )

    # 저장
    pptx_path = generator.save(output_path)

    # PDF 내보내기 (설정에서 요청한 경우)
    if settings.get("export_pdf", False) or export_pdf:
        pdf_path = generator.export_pdf(pptx_path=str(pptx_path))
        if pdf_path:
            print(f"PDF 생성 완료: {pdf_path}")

    return pptx_path


def main():
    parser = argparse.ArgumentParser(
        description="템플릿 기반 PowerPoint 생성기",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
사용 예시:
  python ppt_generator.py -c config.yaml -o output.pptx
  python ppt_generator.py --title "보고서" --author "홍길동" --type 보고

슬라이드 타입:
  - section: 섹션 구분 슬라이드
  - content: 기본 내용 슬라이드
  - content_icons: 아이콘이 있는 내용
  - comparison: 좌우 비교 슬라이드
  - text: 자유 형식 텍스트
  - table: 표
  - cards: 카드 그리드
  - architecture: 아키텍처 다이어그램
  - flowchart: 플로우차트
  - summary: 요약 슬라이드
  - image: 이미지 슬라이드
  - timeline: 타임라인/로드맵 슬라이드
  - stats: 통계/수치 강조 슬라이드
  - two_column: 2단 레이아웃 슬라이드
  - chart: 차트 슬라이드 (bar, column, line, pie 등)
  - org_chart: 조직도 슬라이드

설정 옵션 (YAML/JSON):
  settings:
    show_page_numbers: true  # 페이지 번호 표시 여부
    export_pdf: false        # PDF 내보내기 여부
    theme_name: default      # 테마 (default, dark, green, purple, warm)
    theme_path: theme.yaml   # 외부 테마 파일 경로
"""
    )

    parser.add_argument("-c", "--config", help="설정 파일 경로 (JSON/YAML)")
    parser.add_argument("-o", "--output", help="출력 파일 경로")
    parser.add_argument("--pdf", action="store_true", help="PDF로도 내보내기")
    parser.add_argument("--theme", choices=["default", "dark", "green", "purple", "warm"],
                        help="테마 선택")
    parser.add_argument("--theme-file", help="외부 테마 파일 경로")
    parser.add_argument("--save-theme", help="현재 테마를 파일로 저장")
    parser.add_argument("--title", help="표지 제목 (간단 모드)")
    parser.add_argument("--date", help="날짜 (간단 모드)")
    parser.add_argument("--author", help="작성자 (간단 모드)")
    parser.add_argument("--type", choices=["의사결정", "보고", "정보공유"],
                        default="정보공유", help="보고 유형 (간단 모드)")

    args = parser.parse_args()

    # 테마 저장 모드
    if args.save_theme:
        design = DesignSystem(theme_name=args.theme, theme_path=args.theme_file)
        design.save_theme(args.save_theme)
        print(f"테마 저장 완료: {args.save_theme}")
        sys.exit(0)

    if args.config:
        config_path = Path(args.config)

        if config_path.suffix in [".yaml", ".yml"]:
            with open(config_path, "r", encoding="utf-8") as f:
                config = yaml.safe_load(f)
        else:
            with open(config_path, "r", encoding="utf-8") as f:
                config = json.load(f)

        # CLI 테마 옵션을 config에 병합
        if args.theme:
            config.setdefault("settings", {})["theme_name"] = args.theme
        if args.theme_file:
            config.setdefault("settings", {})["theme_path"] = args.theme_file

        # 출력 경로 처리 - 항상 output 폴더에 저장
        final_output = args.output
        if final_output:
            output_file = Path(final_output)
            # 경로가 파일명만 있으면 output 폴더에 저장
            if output_file.parent == Path(".") or str(output_file.parent) == ".":
                OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
                final_output = str(OUTPUT_DIR / output_file.name)

        output_path = create_from_config(config, final_output, export_pdf=args.pdf)
        print(f"PPT 생성 완료: {output_path}")

    elif args.title:
        generator = PPTGenerator(theme_name=args.theme, theme_path=args.theme_file)
        generator.add_cover_slide(
            title=args.title,
            date=args.date,
            author=args.author or "",
            report_type=args.type
        )
        # 출력 경로 처리 - 항상 output 폴더에 저장
        final_output = args.output
        if final_output:
            output_file = Path(final_output)
            if output_file.parent == Path(".") or str(output_file.parent) == ".":
                OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
                final_output = str(OUTPUT_DIR / output_file.name)
        output_path = generator.save(final_output)
        print(f"PPT 생성 완료: {output_path}")

        if args.pdf:
            pdf_path = generator.export_pdf(pptx_path=str(output_path))
            if pdf_path:
                print(f"PDF 생성 완료: {pdf_path}")

    else:
        parser.print_help()
        sys.exit(1)


if __name__ == "__main__":
    main()
